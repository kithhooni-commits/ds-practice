"""다른 pptx 의 슬라이드를 발표 파일 **뒤에 붙인다.**

`patch_ppt.py` 는 있는 것을 갈아 끼울 뿐 새 장을 못 만든다. 한 장짜리로 뽑은
그림 슬라이드를 손으로 복사해 넣는 대신 여기서 붙인다.

python-pptx 에 슬라이드 복사 기능이 없어 XML 을 통째로 옮긴다. 그림·글꼴 같은
딸린 부품(관계, relationship)은 새 파일 쪽에 다시 등록해야 하므로, 원본의
관계를 훑어 하나씩 옮겨 붙이고 참조 id 를 새 것으로 바꾼다.

    python append_slides.py --file 발표.pptx --add 한장.pptx [한장2.pptx ...]
    python append_slides.py --file 발표.pptx --add 한장.pptx --at 21

`--at` 을 주면 그 번호 **뒤에** 끼워 넣는다 (생략하면 맨 끝).
"""

from __future__ import annotations

import argparse
import copy
import sys
from io import BytesIO
from pathlib import Path

from pptx import Presentation


def copy_slide(src_slide, dst_prs):
    """슬라이드 하나를 dst_prs 맨 뒤에 복제한다."""
    # 레이아웃은 빈 것을 쓴다. 원본도 빈 레이아웃에 직접 그린 것이라 문제없다
    layout = dst_prs.slide_layouts[6]
    new = dst_prs.slides.add_slide(layout)

    # add_slide 가 레이아웃에서 물려준 자리표시자를 지운다
    for shp in list(new.shapes):
        shp._element.getparent().remove(shp._element)

    for shp in src_slide.shapes:
        new.shapes._spTree.append(copy.deepcopy(shp._element))

    # 그림 등 딸린 부품을 옮기고 참조 id 를 새로 매긴다.
    #
    # 원본 쪽 부품 객체를 그대로 relate_to 하면 이름이 같은 부품이 zip 에 두 번
    # 들어가 파일이 깨진다 (Duplicate name 경고). 그림은 **바이트만 꺼내** 이쪽
    # 패키지에 새로 등록하고, 레이아웃은 add_slide 가 이미 붙였으니 건너뛴다.
    remap = {}
    for rid, rel in src_slide.part.rels.items():
        if rel.reltype.endswith("/slideLayout"):
            continue
        if rel.is_external:
            remap[rid] = new.part.rels.get_or_add_ext_rel(rel.reltype, rel.target_ref)
        elif rel.reltype.endswith("/image"):
            _, new_rid = new.part.get_or_add_image_part(BytesIO(rel.target_part.blob))
            remap[rid] = new_rid
        else:
            remap[rid] = new.part.relate_to(rel.target_part, rel.reltype)

    ns = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"
    for el in new.shapes._spTree.iter():
        for attr in ("embed", "link", "id"):
            key = ns + attr
            if key in el.attrib and el.attrib[key] in remap:
                el.attrib[key] = remap[el.attrib[key]]

    # 배경색도 따라가야 한다 (bg() 로 칠한 것)
    bg = src_slide._element.find(
        "{http://schemas.openxmlformats.org/drawingml/2006/main}bg")
    if bg is None:
        bg = src_slide._element.find(
            "{http://schemas.openxmlformats.org/presentationml/2006/main}bg")
    if bg is not None:
        new._element.insert(0, copy.deepcopy(bg))
    return new


def move_to(prs, index_from: int, index_to: int) -> None:
    """맨 뒤에 붙은 슬라이드를 원하는 자리로 옮긴다."""
    lst = prs.slides._sldIdLst
    ids = list(lst)
    lst.remove(ids[index_from])
    lst.insert(index_to, ids[index_from])


def main() -> None:
    try:
        sys.stdout.reconfigure(encoding="utf-8")
    except Exception:
        pass

    ap = argparse.ArgumentParser()
    ap.add_argument("--file", type=Path, required=True)
    ap.add_argument("--add", nargs="+", type=Path, required=True,
                    help="붙일 pptx 들. 그 파일의 모든 슬라이드를 가져온다")
    ap.add_argument("--at", type=int, default=None,
                    help="이 번호 뒤에 끼워 넣는다 (생략하면 맨 끝)")
    ap.add_argument("--out", type=Path, default=None)
    args = ap.parse_args()

    prs = Presentation(str(args.file))
    n0 = len(prs.slides._sldIdLst)

    added = 0
    for path in args.add:
        src = Presentation(str(path))
        for sl in src.slides:
            copy_slide(sl, prs)
            added += 1
        print(f"  {path.name}: {len(src.slides._sldIdLst)}장 가져옴")

    if args.at is not None:
        # 뒤에 쌓인 것들을 순서 유지하며 원하는 자리로 옮긴다
        for k in range(added):
            move_to(prs, n0 + k, args.at + k)
        print(f"  {args.at}번 뒤에 끼워 넣음")

    out = args.out or args.file
    prs.save(str(out))
    print(f"\n{n0}장 → {len(prs.slides._sldIdLst)}장   저장 → {out}")


if __name__ == "__main__":
    main()
