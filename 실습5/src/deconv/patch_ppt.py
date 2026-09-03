"""이미 손으로 고친 발표 파일의 **숫자와 그림만** 바꾼다.

`make_ppt3.py` 는 처음부터 다시 그리므로 손으로 넣은 수정이 전부 날아간다.
결과가 갱신될 때마다 다시 뽑을 수는 없으니, 여기서는 파일을 열어 **텍스트 조각과
그림 파일만 갈아 끼운다.** 나머지는 손대지 않는다.

## 숫자

`--set 옛값=새값` 을 여러 번 준다. 텍스트 상자와 표의 **run 단위**로 바꾼다 —
문단을 통째로 다시 쓰면 부분 굵게 같은 서식이 날아가기 때문이다.

**순서가 중요하다.** 예를 들어 최종 PSNR 30.37 -> 30.44 와 gaussian 30.45 -> 30.37
을 같이 할 때, 앞엣것을 먼저 해야 새로 들어간 30.37 이 다시 바뀌지 않는다.
준 순서대로 적용한다.

## 그림

`--img 슬라이드번호=파일` 로 준다. 그 슬라이드의 첫 그림을 그 파일로 갈아 끼운다.
자리와 크기는 그대로 두고 **이미지 데이터만** 바꾼다.
"""

from __future__ import annotations

import argparse
import sys
from pathlib import Path

from pptx import Presentation


def iter_runs(slide):
    """텍스트 상자와 표 안의 모든 run 을 돌려준다."""
    for sh in slide.shapes:
        if sh.has_text_frame:
            for para in sh.text_frame.paragraphs:
                yield from para.runs
        if sh.has_table:
            for row in sh.table.rows:
                for cell in row.cells:
                    for para in cell.text_frame.paragraphs:
                        yield from para.runs


def main() -> None:
    try:
        sys.stdout.reconfigure(encoding="utf-8")
    except Exception:
        pass

    ap = argparse.ArgumentParser()
    ap.add_argument("--file", type=Path, required=True)
    ap.add_argument("--out", type=Path, default=None, help="생략하면 제자리에 덮어쓴다")
    ap.add_argument("--set", nargs="*", default=[], metavar="옛값=새값",
                    help="준 순서대로 적용한다")
    ap.add_argument("--img", nargs="*", default=[], metavar="슬라이드번호=파일",
                    help="그 슬라이드의 첫 그림을 갈아 끼운다 (자리·크기는 유지)")
    args = ap.parse_args()

    prs = Presentation(str(args.file))

    # ---- 숫자 ----
    pairs = []
    for s in args.set:
        old, _, new = s.partition("=")
        if not _:
            raise SystemExit(f"'옛값=새값' 형식이어야 한다: {s}")
        pairs.append((old, new))

    counts = {o: 0 for o, _ in pairs}
    where: dict[str, list[int]] = {o: [] for o, _ in pairs}
    for i, sl in enumerate(prs.slides, 1):
        for run in iter_runs(sl):
            for old, new in pairs:
                if old in run.text:
                    counts[old] += run.text.count(old)
                    run.text = run.text.replace(old, new)
                    where[old].append(i)

    if pairs:
        print(f"{'옛값':>12} → {'새값':<12}{'바꾼 곳':>8}   슬라이드")
        print("-" * 62)
        for old, new in pairs:
            pg = sorted(set(where[old]))
            mark = "" if counts[old] else "   ← 못 찾음"
            print(f"{old:>12} → {new:<12}{counts[old]:>8}   {pg}{mark}")

    # ---- 그림 ----
    if args.img:
        print()
        for spec in args.img:
            num, _, path = spec.partition("=")
            f = Path(path)
            if not f.exists():
                print(f"  {num}번: {f} 없음 — 건너뜀")
                continue
            sl = list(prs.slides)[int(num) - 1]
            pics = [sh for sh in sl.shapes if sh.shape_type == 13]
            if not pics:
                print(f"  {num}번: 그림이 없다 — 건너뜀")
                continue
            rid = pics[0]._element.blipFill.blip.rEmbed
            part = sl.part.related_part(rid)
            part._blob = f.read_bytes()
            print(f"  {num}번 그림 ← {f.name}  ({f.stat().st_size/1e6:.2f} MB)")

    out = args.out or args.file
    prs.save(str(out))
    print(f"\n저장 → {out}")


if __name__ == "__main__":
    main()
