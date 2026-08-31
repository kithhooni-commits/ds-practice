"""발표 pptx 를 웹에서 훑어볼 수 있는 HTML 로 뽑는다.

내용을 다시 적지 않고 **pptx 에서 그대로 추출**한다. 그래야 슬라이드를 고쳤을 때
미리보기가 조용히 어긋나는 일이 없다. PowerPoint 없이, 폰에서도 검토할 수 있다.

그림은 원본이 장당 1MB 가 넘어 그대로 넣으면 페이지가 10MB 를 넘는다.
가로 1500px, JPEG 82 로 줄여 담는다 — 검토용이므로 충분하다.
"""

from __future__ import annotations

import argparse
import base64
import html
import io
import sys
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Emu

HERE = Path(__file__).resolve().parent
ROOT = HERE.parent.parent


def shrink(blob: bytes, max_w: int = 1500, quality: int = 82) -> str:
    im = Image.open(io.BytesIO(blob))
    if im.mode not in ("RGB", "L"):
        im = im.convert("RGB")
    if im.width > max_w:
        im = im.resize((max_w, round(im.height * max_w / im.width)), Image.LANCZOS)
    buf = io.BytesIO()
    im.save(buf, "JPEG", quality=quality, optimize=True)
    return "data:image/jpeg;base64," + base64.b64encode(buf.getvalue()).decode()


def collect(slide) -> dict:
    """도형을 위→아래 순서로 훑어 제목·본문·표·그림으로 나눈다."""
    items = []
    for sh in sorted(slide.shapes, key=lambda s: (s.top or 0, s.left or 0)):
        if sh.has_text_frame and sh.text_frame.text.strip():
            runs = []
            for p in sh.text_frame.paragraphs:
                t = "".join(r.text for r in p.runs).strip()
                if not t:
                    continue
                size = max((r.font.size.pt for r in p.runs if r.font.size), default=12)
                runs.append((t, size))
            if runs:
                items.append(("text", runs, Emu(sh.top or 0).inches))
        elif sh.has_table:
            rows = [[c.text.strip() for c in r.cells] for r in sh.table.rows]
            items.append(("table", rows, Emu(sh.top or 0).inches))
        elif sh.shape_type == 13:  # PICTURE
            items.append(("image", sh.image.blob, Emu(sh.top or 0).inches))
    return items


def render(items: list) -> tuple[str, str]:
    """(제목, 본문 HTML). 가장 큰 글씨를 제목으로 본다."""
    title = ""
    best = 0.0
    for kind, payload, _ in items:
        if kind != "text":
            continue
        for t, size in payload:
            if size > best:
                best, title = size, t

    out = []
    for kind, payload, _ in items:
        if kind == "text":
            for t, size in payload:
                if t == title:
                    continue
                cls = "lead" if size >= 15 else ("small" if size <= 11 else "body")
                if size <= 11.5 and t.isupper() is False and len(t) < 40 and size < 12:
                    cls = "small"
                out.append(f'<p class="{cls}">{html.escape(t)}</p>')
        elif kind == "table":
            rows = payload
            th = "".join(f"<th>{html.escape(c)}</th>" for c in rows[0])
            body = "".join(
                "<tr>" + "".join(f"<td>{html.escape(c).replace(chr(10), '<br>')}</td>" for c in r) + "</tr>"
                for r in rows[1:]
            )
            out.append(f'<div class="tw"><table><thead><tr>{th}</tr></thead><tbody>{body}</tbody></table></div>')
        elif kind == "image":
            out.append(f'<img src="{shrink(payload)}" alt="">')
    return title, "\n".join(out)


CSS = """
:root{--ground:#F3F6F5;--surface:#FFF;--surface-2:#ECF0EE;--ink:#131A19;--ink-2:#3C4947;
--muted:#66756F;--rule:#DCE3E0;--accent:#0B6E5E;--accent-soft:#DCEDE8;}
@media (prefers-color-scheme:dark){:root:not([data-theme="light"]){--ground:#0F1413;--surface:#161D1C;
--surface-2:#1D2625;--ink:#E8EEEC;--ink-2:#C0CCC9;--muted:#8B9A96;--rule:#28322F;
--accent:#55C4AA;--accent-soft:#17302B;}}
:root[data-theme="dark"]{--ground:#0F1413;--surface:#161D1C;--surface-2:#1D2625;--ink:#E8EEEC;
--ink-2:#C0CCC9;--muted:#8B9A96;--rule:#28322F;--accent:#55C4AA;--accent-soft:#17302B;}
*{box-sizing:border-box}
body{margin:0;background:var(--ground);color:var(--ink);
font-family:"IBM Plex Sans","Malgun Gothic",system-ui,sans-serif;line-height:1.6}
.wrap{max-width:1000px;margin:0 auto;padding:clamp(1.5rem,4vw,3rem) clamp(1rem,3vw,2rem) 5rem;
display:flex;flex-direction:column;gap:1.5rem}
header{border-bottom:1px solid var(--rule);padding-bottom:1.25rem}
h1{font-family:"IBM Plex Sans Condensed","Malgun Gothic",sans-serif;font-size:clamp(1.7rem,4vw,2.4rem);
margin:.3rem 0 .5rem;letter-spacing:-.02em;text-wrap:balance}
.eyebrow{font-family:"IBM Plex Mono",monospace;font-size:.75rem;letter-spacing:.14em;
text-transform:uppercase;color:var(--accent);margin:0;font-weight:600}
.note{color:var(--muted);font-size:.9rem;margin:.4rem 0 0;max-width:62ch}
.slide{background:var(--surface);border:1px solid var(--rule);border-radius:4px;overflow:hidden}
.head{display:flex;align-items:baseline;gap:.7rem;padding:.85rem 1.25rem;
background:var(--surface-2);border-bottom:1px solid var(--rule)}
.num{font-family:"IBM Plex Mono",monospace;font-size:.8rem;color:var(--accent);font-weight:600;
min-width:1.9rem}
.head h2{font-family:"IBM Plex Sans Condensed","Malgun Gothic",sans-serif;font-size:1.15rem;
margin:0;letter-spacing:-.01em}
.body{padding:1.1rem 1.25rem 1.4rem;display:flex;flex-direction:column;gap:.7rem}
p{margin:0;max-width:70ch}
.lead{font-size:1.02rem;color:var(--ink);font-weight:500}
.body p.body{color:var(--ink-2);font-size:.92rem}
.small{color:var(--muted);font-size:.82rem;font-family:"IBM Plex Mono",monospace}
img{max-width:100%;height:auto;border:1px solid var(--rule);border-radius:3px;display:block}
.tw{overflow-x:auto;border:1px solid var(--rule);border-radius:3px}
table{border-collapse:collapse;width:100%;min-width:520px;font-size:.85rem}
th{background:var(--accent-soft);text-align:left;padding:.45rem .7rem;font-weight:600;
border-bottom:1px solid var(--rule);white-space:nowrap}
td{padding:.42rem .7rem;border-bottom:1px solid var(--rule);color:var(--ink-2);vertical-align:top}
td:not(:first-child){font-family:"IBM Plex Mono",monospace;font-variant-numeric:tabular-nums}
tr:last-child td{border-bottom:0}
"""


def main() -> None:
    try:
        sys.stdout.reconfigure(encoding="utf-8")
    except Exception:
        pass

    ap = argparse.ArgumentParser()
    ap.add_argument("--pptx", type=Path, default=ROOT / "실습5_denoising_발표.pptx")
    ap.add_argument("--out", type=Path, required=True)
    args = ap.parse_args()

    prs = Presentation(str(args.pptx))
    cards = []
    for i, slide in enumerate(prs.slides, 1):
        title, body = render(collect(slide))
        cards.append(
            f'<section class="slide"><div class="head"><span class="num">{i:02d}</span>'
            f"<h2>{html.escape(title)}</h2></div>"
            f'<div class="body">{body}</div></section>'
        )

    doc = (
        "<title>Day 1 발표 미리보기</title>\n"
        '<link rel="preconnect" href="https://fonts.googleapis.com">\n'
        '<link rel="preconnect" href="https://fonts.gstatic.com" crossorigin>\n'
        '<link rel="stylesheet" href="https://fonts.googleapis.com/css2?'
        "family=IBM+Plex+Mono:wght@400;600&family=IBM+Plex+Sans+Condensed:wght@600&"
        'family=IBM+Plex+Sans:wght@400;500;600&display=swap">\n'
        f"<style>{CSS}</style>\n"
        '<div class="wrap"><header>'
        '<p class="eyebrow">삼성 DS2 · Image Restoration Challenge · Day 1</p>'
        f"<h1>발표 자료 미리보기 — {len(cards)}장</h1>"
        '<p class="note">pptx 에서 그대로 추출했다. 내용 검토용이라 실제 슬라이드의 배치·색은 '
        "재현하지 않는다. 그림은 검토용으로 줄여 담았다.</p>"
        "</header>\n" + "\n".join(cards) + "</div>"
    )

    args.out.write_text(doc, encoding="utf-8")
    print(f"{len(cards)}장 → {args.out}  ({args.out.stat().st_size / 1e6:.1f} MB)")


if __name__ == "__main__":
    main()
