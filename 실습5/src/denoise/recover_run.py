"""발표 pptx 에서 학습 결과를 되살린다.

Colab 런타임이 끊기면 `/content/runs` 가 통째로 사라진다. 체크포인트는 못 살리지만,
**슬라이드에 이미 박혀 있는 수치와 그림은 살릴 수 있다.** 그것만 있으면 슬라이드를
다시 만들 수 있으므로, pptx 를 역으로 읽어 `test_metrics*.json` 과 `figures/*.png`
를 복원한다.

되살리는 것
  - 노이즈 종류별 model PSNR (슬라이드 5 표)
  - 전체 PSNR / SSIM, salt&pepper SSIM (슬라이드 5 본문)
  - self-ensemble 기여도 (슬라이드 4 표) → self-ensemble 없는 쪽 PSNR 을 역산
  - 4패널·오차지도·막대그래프 그림 (슬라이드에 박힌 원본 해상도 PNG)

되살리지 못하는 것
  - 체크포인트. 다시 추론할 수는 없다
  - gaussian/rician/uniform 의 SSIM. 슬라이드에 안 나오는 값이라 전체 SSIM 이
    맞도록 배분한다. **측정값이 아니므로 표시하지 않는다** (실제로 어떤 슬라이드에도
    안 쓰인다). JSON 의 `estimated` 필드로 표시해 둔다.

conventional 필터(noisy/mean/median/adaptive) 수치는 모델과 무관하므로 로컬에
있는 아무 평가 결과에서 그대로 가져온다.
"""

from __future__ import annotations

import argparse
import json
import re
import sys
from pathlib import Path

from pptx import Presentation

HERE = Path(__file__).resolve().parent
ROOT = HERE.parent.parent
NOISE_ORDER = ["gaussian", "rician", "uniform", "salt_and_pepper"]
LABEL2KEY = {"Gaussian": "gaussian", "Rician": "rician",
             "Uniform": "uniform", "Salt & Pepper": "salt_and_pepper"}
# 슬라이드 6-9 는 NOISE_ORDER 순서, 그 다음이 오차지도
PANEL_SLIDES = {7: "panel_gaussian", 8: "panel_rician", 9: "panel_uniform",
                10: "panel_salt_and_pepper", 11: "error_maps", 5: "summary_bars"}


def tables(slide):
    return [sh.table for sh in slide.shapes if sh.has_table]


def main() -> None:
    try:
        sys.stdout.reconfigure(encoding="utf-8")
    except Exception:
        pass

    ap = argparse.ArgumentParser()
    ap.add_argument("pptx", type=Path)
    ap.add_argument("--out", type=Path, required=True, help="복원할 run 폴더")
    ap.add_argument("--template", type=Path, required=True,
                    help="conventional 수치를 가져올 기존 test_metrics_se.json")
    ap.add_argument("--figures", type=Path, default=ROOT / "figures")
    args = ap.parse_args()

    prs = Presentation(str(args.pptx))
    slides = list(prs.slides)

    # ---- 슬라이드 5: 노이즈별 PSNR ----
    per_noise: dict[str, float] = {}
    total_psnr = None
    for row in tables(slides[4])[0].rows:
        cells = [c.text.strip() for c in row.cells]
        if cells[0] in LABEL2KEY:
            per_noise[LABEL2KEY[cells[0]]] = float(cells[2])
        elif cells[0] == "전체":
            total_psnr = float(cells[2])
    assert len(per_noise) == 4 and total_psnr, "슬라이드 5 표를 못 읽었다"

    # ---- 슬라이드 5 본문: 전체 SSIM, s&p SSIM ----
    body = " ".join(sh.text_frame.text for sh in slides[4].shapes if sh.has_text_frame)
    total_ssim = float(re.search(r"SSIM\s+([\d.]+)", body).group(1))
    snp_ssim = float(re.search(r"SSIM\s+[\d.]+\s*→\s*([\d.]+)", body).group(1))

    # ---- 슬라이드 4: self-ensemble 기여도 → 비-SE PSNR 역산 ----
    se_gain = None
    for row in tables(slides[3])[0].rows:
        joined = " ".join(c.text for c in row.cells)
        if "self-ensemble" in joined:
            m = re.search(r"\+([\d.]+)\s*dB", joined)
            if m:
                se_gain = float(m.group(1))
    print(f"읽은 값: PSNR {total_psnr} / SSIM {total_ssim} / s&p SSIM {snp_ssim} / SE +{se_gain}")

    # ---- SSIM 배분 (전체가 맞도록. 슬라이드에 안 쓰이는 값이다) ----
    tmpl = json.loads(args.template.read_text(encoding="utf-8"))
    base_ssim = {nz: sum(r["ssim_model"] for r in tmpl["rows"] if r["noise_type"] == nz)
                 / sum(1 for r in tmpl["rows"] if r["noise_type"] == nz)
                 for nz in NOISE_ORDER}
    rest = [nz for nz in NOISE_ORDER if nz != "salt_and_pepper"]
    want = total_ssim * 4 - snp_ssim
    scale = want / sum(base_ssim[nz] for nz in rest)
    ssim_by = {nz: base_ssim[nz] * scale for nz in rest} | {"salt_and_pepper": snp_ssim}

    # ---- rows 재구성: conventional 은 그대로, model 만 갈아끼운다 ----
    def build_rows(psnr_shift: float = 0.0):
        out = []
        for r in tmpl["rows"]:
            nz = r["noise_type"]
            out.append(r | {"psnr_model": per_noise[nz] + psnr_shift,
                            "ssim_model": ssim_by[nz]})
        return out

    args.out.mkdir(parents=True, exist_ok=True)
    (args.out / "config.json").write_text(
        json.dumps({"model": "drunet", "features": 64, "epochs": 180,
                    "note": "pptx 에서 복원한 run. 체크포인트는 없다."},
                   ensure_ascii=False, indent=2), encoding="utf-8")

    (args.out / "test_metrics_se.json").write_text(json.dumps(
        {"psnr_total": total_psnr, "ssim_total": total_ssim, "self_ensemble": True,
         "recovered_from": str(args.pptx),
         "estimated": ["ssim_model (gaussian/rician/uniform) — 슬라이드에 없어 전체값에 맞춰 배분"],
         "rows": build_rows()}, ensure_ascii=False, indent=2), encoding="utf-8")

    if se_gain:
        (args.out / "test_metrics.json").write_text(json.dumps(
            {"psnr_total": total_psnr - se_gain, "ssim_total": total_ssim, "self_ensemble": False,
             "recovered_from": str(args.pptx),
             "estimated": ["self-ensemble 기여도에서 역산한 값"],
             "rows": build_rows(-se_gain)}, ensure_ascii=False, indent=2), encoding="utf-8")

    # ---- 그림 ----
    args.figures.mkdir(parents=True, exist_ok=True)
    saved = 0
    for idx, name in PANEL_SLIDES.items():
        for sh in slides[idx - 1].shapes:
            if sh.shape_type == 13:
                (args.figures / f"{name}.png").write_bytes(sh.image.blob)
                saved += 1
                break

    print(f"복원: {args.out}")
    print(f"      test_metrics_se.json (PSNR {total_psnr:.2f} / SSIM {total_ssim:.4f})")
    if se_gain:
        print(f"      test_metrics.json    (PSNR {total_psnr - se_gain:.2f}, SE 기여 역산)")
    print(f"      그림 {saved}개 → {args.figures}")


if __name__ == "__main__":
    main()
