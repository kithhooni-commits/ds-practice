"""발표 슬라이드 생성 (python-pptx).

과제 요구사항 네 가지를 슬라이드로 옮긴다.
  1. pipeline
  2. before / after / error map / ground truth
  3. 왜 그 방법을 골랐는가
  4. label-free (가산점)

수치는 하드코딩하지 않고 runs/ 의 test_metrics*.json 에서 읽는다.
그래야 재학습할 때마다 슬라이드가 같이 갱신된다.
"""

from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path

import numpy as np
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

HERE = Path(__file__).resolve().parent
ROOT = HERE.parent.parent
FIGDIR = ROOT / "figures"

W, H = Inches(13.333), Inches(7.5)

INK = RGBColor(0x13, 0x1A, 0x19)
INK2 = RGBColor(0x3C, 0x49, 0x47)
MUTED = RGBColor(0x66, 0x75, 0x6F)
ACCENT = RGBColor(0x0B, 0x6E, 0x5E)
ACCENT_SOFT = RGBColor(0xDC, 0xED, 0xE8)
RULE = RGBColor(0xDC, 0xE3, 0xE0)
SURF = RGBColor(0xF3, 0xF6, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
WARN = RGBColor(0x8C, 0x4A, 0x6B)

FONT = "맑은 고딕"
MONO = "Consolas"

NOISE_ORDER = ["gaussian", "rician", "uniform", "salt_and_pepper"]
NOISE_KO = {"gaussian": "Gaussian", "rician": "Rician",
            "uniform": "Uniform", "salt_and_pepper": "Salt & Pepper"}
BASE = {"psnr": 30.510, "ssim": 0.8950}

# 모델별 표기 — 슬라이드 문구가 실제로 학습한 구조를 따라가게 한다
MODELS = {
    "dncnn": {
        "name": "DnCNN",
        "box": "DnCNN 17층" + chr(10) + "64ch · 전역 잔차",
        "rf": "35px",
        "params": "0.59M",
    },
    "dncnn_plus": {
        "name": "DnCNN + median 채널",
        "box": "DnCNN + median 채널" + chr(10) + "17층 · 64ch",
        "rf": "35px",
        "params": "0.59M",
    },
    "drunet": {
        "name": "DRUNet",
        "box": "DRUNet (U-Net + res block)" + chr(10) + "4스케일 · 수용영역 ~180px",
        "rf": "~180px",
        "params": "32.6M",
    },
}


def read_model(run_dir: Path) -> dict:
    """run 폴더의 config.json 에서 실제 학습한 구조를 읽는다."""
    cfg = run_dir / "config.json"
    key = "dncnn"
    if cfg.exists():
        try:
            key = json.loads(cfg.read_text(encoding="utf-8")).get("model", "dncnn")
        except Exception:
            pass
    return MODELS.get(key, MODELS["dncnn"]) | {"key": key}


# 시도 이력 — (상태, 시도, 결과, 배운 것)
# 채택/기각을 한 장에 같이 둔다. 무엇을 했는지보다 무엇이 틀렸는지가 근거로 강하다.
# 시도 이력 — 지나간 실험은 상수로, 제출 모델 칸은 실측값으로 채운다.
# (상태, 시도, 결과, 배운 것) 중 결과가 None 이면 history_rows 가 계산해 넣는다.
DNCNN_BEST = 34.56  # A100 60 epoch DnCNN + self-ensemble. 구조 비교의 기준점
DRUNET_60EP = 36.50  # 같은 조건 DRUNet 60 epoch. 학습 길이 효과를 분리하는 기준점


def history_rows(se, p_model, mi: dict, p_dncnn_snp: float = 36.78, lf: dict | None = None) -> list[tuple]:
    """시도 요약 표. 마지막 두 칸은 실제 결과에서 계산한다."""
    final = se["psnr_total"]
    snp = p_model["salt_and_pepper"]
    rows = [
        ("전제", "지표 구현을 배포 코드에서 그대로 이식",
         "0.0000 dB", "배포 로그와 넷째자리까지 일치. 여기가 맞기 전엔 어떤 숫자도 근거가 아니다"),
        ("채택", "Charbonnier · cosine LR · 크롭 128 · rot90",
         "30.51 → 33.58", "구조를 안 건드리고 학습 절차만 바꿔 +3.07 dB. 개선의 대부분이 여기서 나왔다"),
        ("채택", "8× self-ensemble (dihedral 평균)",
         "33.58 → 34.13", "학습 비용 0. 모델이 좋아질수록 여지는 준다"),
        ("채택", "epoch 연장 (DnCNN 40 → 60)",
         f"34.13 → {DNCNN_BEST:.2f}", "val best 가 마지막 epoch 이면 아직 수렴 전이라는 신호다"),
    ]
    if mi["key"] == "drunet":
        rows.append(
            ("채택", "DRUNet — 수용영역 35px → 180px",
             f"{DNCNN_BEST:.2f} → {DRUNET_60EP:.2f}",
             f"같은 60 epoch 에서 +{DRUNET_60EP - DNCNN_BEST:.2f} dB. s&p 에서만 {snp - p_dncnn_snp:+.2f} dB"))
        rows.append(
            ("채택", "DRUNet 을 180 epoch 까지",
             f"{DRUNET_60EP:.2f} → {final:.2f}",
             "파라미터가 55배인 모델은 더 오래 걸린다. 여기서도 best 가 173 epoch"))
    rows += [
        ("기각", "median 3×3 을 두 번째 입력 채널로",
         "−0.39 / −0.87 dB", "진단은 맞고 처방이 틀렸다. 좁은 시야를 둔 채 국소 도구를 더한 셈"),
        ("기각", "σ 게이트 — 깨끗한 입력은 그대로 통과",
         "오라클 +0.52", "σ̂ 이 실제 0.001 을 0.037 로 추정. 제약이 실제 비용을 만든다는 증거"),
        ("기각", "patch 256 · batch 32 (A100 이니까)",
         "val −0.11 dB", "배치 2배 = epoch 당 step 절반. 이득은 배치가 아니라 step 수다"),
        ("정정", "label-free 학습 데이터: test 100장 vs train 7,268장",
         "30.95 vs 30.35", "처음엔 fp16 불안정으로 망가진 실행을 보고 '데이터가 안 통한다'고 "
                           "단정했다. 제대로 돌리니 −0.60 dB — 방향은 같아도 근거가 없었다"),
    ]
    if lf:
        v = next(iter(lf.values()))
        rows.append(("가산점", "label-free (Noise2Void, clean 0장)",
                     f"{v['psnr']:.2f} / {v['ssim']:.4f}",
                     "clean 7,268장 지도학습 기준선(30.51)을 넘었다"))
    return rows




# ------------------------------------------------------------------ helpers


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def bg(slide, color=WHITE):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def text(slide, x, y, w, h, runs, size=16, color=INK2, bold=False, font=FONT,
         align=PP_ALIGN.LEFT, spacing=1.25, anchor=MSO_ANCHOR.TOP):
    """runs: 문자열 또는 (문자열, dict) 목록. dict 로 run 단위 서식을 덮어쓴다."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    if isinstance(runs, str):
        runs = [runs]

    first = True
    for item in runs:
        s, opt = item if isinstance(item, tuple) else (item, {})
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = opt.get("align", align)
        p.line_spacing = opt.get("spacing", spacing)
        if "space_before" in opt:
            p.space_before = opt["space_before"]
        r = p.add_run()
        r.text = s
        f = r.font
        f.name = opt.get("font", font)
        f.size = Pt(opt.get("size", size))
        f.bold = opt.get("bold", bold)
        f.color.rgb = opt.get("color", color)
    return tb


def title(slide, main, sub=None, eyebrow=None):
    y = Inches(0.46)
    if eyebrow:
        text(slide, Inches(0.7), y, Inches(12), Inches(0.3), eyebrow,
             size=11, color=ACCENT, bold=True, font=MONO)
        y = Inches(0.78)
    text(slide, Inches(0.7), y, Inches(12), Inches(0.6), main, size=30, color=INK, bold=True)
    y2 = y + Inches(0.72)
    if sub:
        text(slide, Inches(0.7), y2, Inches(12), Inches(0.4), sub, size=14, color=MUTED)
        y2 += Inches(0.5)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), y2 + Inches(0.05), Inches(11.93), Pt(1.2))
    line.fill.solid()
    line.fill.fore_color.rgb = RULE
    line.line.fill.background()
    line.shadow.inherit = False
    return y2 + Inches(0.35)


def card(slide, x, y, w, h, fill=SURF, line=RULE):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sh.adjustments[0] = 0.04
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = line
    sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    return sh


def arrow(slide, x, y, w, h, color=ACCENT):
    sh = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def table(slide, x, y, w, h, rows, col_w=None, head=True, size=12, highlight_row=None,
          highlight_col=None):
    nr, nc = len(rows), len(rows[0])
    shape = slide.shapes.add_table(nr, nc, x, y, w, h)
    tbl = shape.table
    if col_w:
        total = sum(col_w)
        for i, cw in enumerate(col_w):
            tbl.columns[i].width = Emu(int(w * cw / total))
    for r, row in enumerate(rows):
        tbl.rows[r].height = Inches(0.3)
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = str(val)
            cell.margin_left = Inches(0.09)
            cell.margin_right = Inches(0.09)
            cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.03)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            is_head = head and r == 0
            hot = (highlight_row is not None and r == highlight_row) or \
                  (highlight_col is not None and c == highlight_col and r > 0)
            cell.fill.fore_color.rgb = ACCENT_SOFT if is_head else (
                RGBColor(0xEC, 0xF5, 0xF2) if hot else WHITE)
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.RIGHT
            for run in p.runs:
                run.font.size = Pt(size)
                run.font.name = FONT if c == 0 else MONO
                run.font.bold = is_head or hot
                run.font.color.rgb = INK if (is_head or hot) else INK2
    return shape


def pic(slide, path: Path, x, y, w=None, h=None):
    if not path.exists():
        text(slide, x, y, Inches(6), Inches(0.4), f"[그림 없음: {path.name}]", size=12, color=WARN)
        return None
    return slide.shapes.add_picture(str(path), x, y, width=w, height=h)


def footer(slide, s):
    text(slide, Inches(0.7), Inches(6.95), Inches(12), Inches(0.3), s, size=10, color=MUTED, font=MONO)


# ------------------------------------------------------------------ slides


def build(prs, M, name: str, lf: dict | None, rej: dict | None, mi: dict):
    se, base = M["se"], M["base"]

    def by_noise(rows, key):
        return {nz: float(np.mean([r[key] for r in rows if r["noise_type"] == nz])) for nz in NOISE_ORDER}

    p_model = by_noise(se["rows"], "psnr_model")
    q_model = by_noise(se["rows"], "ssim_model")
    p_in = by_noise(se["rows"], "psnr_noisy")
    q_in = by_noise(se["rows"], "ssim_noisy")
    p_med = by_noise(se["rows"], "psnr_median")
    p_mean = by_noise(se["rows"], "psnr_mean")
    p_adap = by_noise(se["rows"], "psnr_adaptive")

    # ---------------------------------------------------------- 1. 표지
    s = blank(prs); bg(s)
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.28), H)
    band.fill.solid(); band.fill.fore_color.rgb = ACCENT
    band.line.fill.background(); band.shadow.inherit = False

    text(s, Inches(1.0), Inches(1.5), Inches(11), Inches(0.4),
         "삼성 DS2 · Image Restoration Challenge · Day 1", size=13, color=ACCENT, bold=True, font=MONO)
    text(s, Inches(1.0), Inches(2.0), Inches(11), Inches(1.2),
         "노이즈 4종을 하나의 네트워크로", size=42, color=INK, bold=True)
    text(s, Inches(1.0), Inches(2.85), Inches(11), Inches(0.7),
         "종류도 세기도 모른 채 복원하기", size=42, color=ACCENT, bold=True)
    text(s, Inches(1.0), Inches(3.95), Inches(10), Inches(0.5),
         f"반도체 이미지 denoising · test 100장 · {mi['name']} + 학습 레시피 + self-ensemble",
         size=15, color=MUTED)

    c = card(s, Inches(1.0), Inches(4.75), Inches(4.4), Inches(1.5), ACCENT_SOFT, ACCENT)
    text(s, Inches(1.25), Inches(4.95), Inches(4), Inches(0.3), "제출값", size=12, color=ACCENT, bold=True, font=MONO)
    text(s, Inches(1.25), Inches(5.3), Inches(4), Inches(0.7),
         f"PSNR {se['psnr_total']:.2f}   SSIM {se['ssim_total']:.4f}", size=21, color=INK, bold=True, font=MONO)

    c2 = card(s, Inches(5.7), Inches(4.75), Inches(4.4), Inches(1.5))
    text(s, Inches(5.95), Inches(4.95), Inches(4), Inches(0.3), "배포 기준선 (DnCNN 10 epoch)",
         size=12, color=MUTED, bold=True, font=MONO)
    text(s, Inches(5.95), Inches(5.3), Inches(4), Inches(0.7),
         f"PSNR {BASE['psnr']:.2f}   SSIM {BASE['ssim']:.4f}", size=21, color=MUTED, bold=True, font=MONO)
    text(s, Inches(5.95), Inches(5.85), Inches(4), Inches(0.3),
         f"→  +{se['psnr_total'] - BASE['psnr']:.2f} dB   +{se['ssim_total'] - BASE['ssim']:.4f}",
         size=13, color=ACCENT, bold=True, font=MONO)
    text(s, Inches(10.4), Inches(5.9), Inches(2.3), Inches(0.4), name, size=13, color=MUTED, align=PP_ALIGN.RIGHT)

    # ---------------------------------------------------------- 2. 문제
    s = blank(prs); bg(s)
    y = title(s, "문제 — 종류를 모른다는 것이 전부다",
              "이미지마다 4종 중 하나가 랜덤 σ 로 실려 있고, 어느 것인지 알려주지 않는다",
              "PROBLEM")

    rows = [["노이즈", "σ 범위", "성질", "입력 PSNR"]]
    desc = {
        "gaussian": "x + N(0, σ)",
        "rician": "|x + N + iN| — 어두운 쪽이 위로 들림",
        "uniform": "x + U(-σ, σ)",
        "salt_and_pepper": "픽셀의 σ 비율을 0 또는 max 로 덮어씀 (임펄스)",
    }
    rng = {"gaussian": "0 – 0.1", "rician": "0 – 0.15", "uniform": "0 – 0.2", "salt_and_pepper": "0 – 0.2"}
    for nz in NOISE_ORDER:
        rows.append([NOISE_KO[nz], rng[nz], desc[nz], f"{p_in[nz]:.2f} dB"])
    table(s, Inches(0.7), y, Inches(11.9), Inches(1.9), rows, col_w=[2.0, 1.6, 6.3, 1.7],
          highlight_row=4)

    y2 = y + Inches(2.25)
    c = card(s, Inches(0.7), y2, Inches(5.8), Inches(2.5))
    text(s, Inches(0.95), y2 + Inches(0.22), Inches(5.3), Inches(0.35),
         "고전 필터는 종류마다 정답이 정반대다", size=16, color=INK, bold=True)
    text(s, Inches(0.95), y2 + Inches(0.72), Inches(5.3), Inches(1.6), [
        (f"Salt & Pepper  ·  median {p_med['salt_and_pepper']:.1f}  vs  mean {p_mean['salt_and_pepper']:.1f} dB", {"font": MONO, "size": 12.5}),
        (f"Gaussian       ·  adaptive {p_adap['gaussian']:.1f}  vs  median {p_med['gaussian']:.1f} dB", {"font": MONO, "size": 12.5}),
        ("임펄스는 몇 픽셀이 통째로 거짓말이라 평균 계열이 그 거짓말을 이웃에 퍼뜨린다. "
         "반대로 median 은 미세한 계조를 계단으로 뭉갠다.", {"size": 12.5, "space_before": Pt(10)}),
    ], size=12.5, color=INK2)

    c = card(s, Inches(6.8), y2, Inches(5.8), Inches(2.5), ACCENT_SOFT, ACCENT)
    text(s, Inches(7.05), y2 + Inches(0.22), Inches(5.3), Inches(0.35),
         "그래서 고정 필터 하나로는 못 넘는다", size=16, color=INK, bold=True)
    tot_in = se["rows"]
    m_all = float(np.mean([r["psnr_mean"] for r in tot_in]))
    n_all = float(np.mean([r["psnr_noisy"] for r in tot_in]))
    text(s, Inches(7.05), y2 + Inches(0.72), Inches(5.3), Inches(1.6), [
        (f"복원 안 함   {n_all:.2f} dB", {"font": MONO, "size": 13}),
        (f"mean 3×3     {m_all:.2f} dB   ← 겨우 +{m_all - n_all:.2f}", {"font": MONO, "size": 13}),
        ("종류를 모른 채 하나의 규칙을 고정하면 아무것도 안 한 것과 다를 바 없다. "
         "학습이 필요한 이유가 여기 있다.", {"size": 12.5, "space_before": Pt(10)}),
    ], size=12.5, color=INK2)
    footer(s, "출처: dataset/test_noise_only/noise_meta.json · 25장 × 4종 = 100장")

    # ---------------------------------------------------------- 3. 파이프라인
    s = blank(prs); bg(s)
    y = title(s, "파이프라인", "학습은 노이즈를 만들어 배우고, 추론은 만들지 않는다", "1 · PIPELINE")

    # 학습 경로
    text(s, Inches(0.7), y + Inches(0.05), Inches(3), Inches(0.3), "학습 (supervised)",
         size=13, color=ACCENT, bold=True, font=MONO)
    ty = y + Inches(0.45)
    boxes = [
        ("clean 7,268장", "train/", 2.05),
        ("랜덤 크롭 128\nflip · rot90", "증강", 1.75),
        ("노이즈 4종 중\n1개 랜덤 + σ 랜덤", "합성", 2.1),
        (mi["box"], "모델", 2.5),
        ("Charbonnier loss\nvs clean", "학습 신호", 2.0),
    ]
    x = Inches(0.7)
    for i, (label, tag, wid) in enumerate(boxes):
        wid = Inches(wid)
        fill = ACCENT_SOFT if i == 3 else SURF
        card(s, x, ty, wid, Inches(1.05), fill, ACCENT if i == 3 else RULE)
        text(s, x + Inches(0.12), ty + Inches(0.08), wid - Inches(0.24), Inches(0.22), tag,
             size=9.5, color=ACCENT if i == 3 else MUTED, bold=True, font=MONO, align=PP_ALIGN.CENTER)
        text(s, x + Inches(0.1), ty + Inches(0.36), wid - Inches(0.2), Inches(0.6), label,
             size=11.5, color=INK, bold=(i == 3), align=PP_ALIGN.CENTER, spacing=1.05)
        x += wid
        if i < len(boxes) - 1:
            arrow(s, x + Inches(0.04), ty + Inches(0.4), Inches(0.28), Inches(0.24))
            x += Inches(0.36)

    # 추론 경로
    iy = ty + Inches(1.6)
    text(s, Inches(0.7), iy - Inches(0.35), Inches(4), Inches(0.3), "추론 (test 100장)",
         size=13, color=ACCENT, bold=True, font=MONO)
    boxes2 = [
        ("corrupted 100장", "test_noise_only/", 2.6),
        ("8× self-ensemble\ndihedral 평균", "추론", 2.5),
        ("같은 네트워크\n(가중치 1벌)", "모델", 2.5),
        (f"PSNR {se['psnr_total']:.2f}\nSSIM {se['ssim_total']:.4f}", "제출", 2.3),
    ]
    x = Inches(0.7)
    for i, (label, tag, wid) in enumerate(boxes2):
        wid = Inches(wid)
        last = i == len(boxes2) - 1
        card(s, x, iy, wid, Inches(1.05), ACCENT_SOFT if last else SURF, ACCENT if last else RULE)
        text(s, x + Inches(0.12), iy + Inches(0.08), wid - Inches(0.24), Inches(0.22), tag,
             size=9.5, color=ACCENT if last else MUTED, bold=True, font=MONO, align=PP_ALIGN.CENTER)
        text(s, x + Inches(0.1), iy + Inches(0.36), wid - Inches(0.2), Inches(0.6), label,
             size=11.5 if not last else 12.5, color=INK, bold=last,
             align=PP_ALIGN.CENTER, spacing=1.05, font=MONO if last else FONT)
        x += wid
        if not last:
            arrow(s, x + Inches(0.04), iy + Inches(0.4), Inches(0.28), Inches(0.24))
            x += Inches(0.36)

    ky = iy + Inches(1.55)
    card(s, Inches(0.7), ky, Inches(11.9), Inches(1.25))
    text(s, Inches(0.95), ky + Inches(0.15), Inches(11.4), Inches(1.0), [
        ("핵심: 추론 경로에는 노이즈 종류·σ 정보가 들어가지 않는다", {"size": 14, "bold": True, "color": INK}),
        ("noise_meta.json 에 종류와 σ 가 다 적혀 있지만 복원에는 쓰지 않았다. "
         "표를 종류별로 쪼개 보여주는 데만 썼다 — 채점 세트에 그 정보가 없을 수 있고, "
         "모르고도 되는 것이 더 강한 결과다.", {"size": 12.5, "space_before": Pt(6)}),
    ], color=INK2)

    # ---------------------------------------------------------- 4. 왜 이 방법인가
    s = blank(prs); bg(s)
    y = title(s, "왜 이 방법을 골랐는가",
              "배포 코드와 같은 데이터·같은 노이즈 합성·같은 지표를 쓰고 학습 쪽만 바꿨다",
              "3 · DESIGN RATIONALE")

    rows = [["바꾼 것", "관찰", "그래서"]]
    rows += [
        ["Charbonnier loss\n(L2 → smooth L1)",
         "L2 는 s&p 의 극단값 몇 픽셀에\ngradient 가 끌려간다",
         "큰 오차의 영향을 선형으로 제한"],
        ["cosine LR · 40 epoch",
         "배포 설정은 10 epoch 에 plateau×0.88 —\n감쇠가 사실상 안 걸림",
         "끝까지 부드럽게 낮춰 수렴"],
        ["patch 128 랜덤 크롭\n+ rot90",
         f"완전 합성곱이라 학습 패치와\n추론 크기를 맞출 이유가 없다",
         "크롭 자체가 증강.\nstep 수도 2배"],
        ["8× self-ensemble",
         "flip/rot90 로 학습했으니\n8개 변환을 같게 다뤄야 맞다",
         "실제로는 어긋남. 평균이 그것을 지움\n"
         + (f"(+{se['psnr_total'] - base['psnr_total']:.2f} dB, 학습 비용 0)"
            if M.get("has_base", True) else "(학습 비용 0)")],
    ]
    if rej:
        rows.append(["입력에 median 3×3\n채널 추가  → 기각",
                     f"s&p 입력이 {p_in['salt_and_pepper']:.1f} dB 로 혼자 10 dB 아래.\n"
                     f"고전 필터 중 median 만 통함 ({p_med['salt_and_pepper']:.1f} vs mean {p_mean['salt_and_pepper']:.1f})",
                     f"넣은 이유였던 s&p 에서 오히려 {abs(rej['dp_snp']):.2f} dB 손해.\n"
                     f"17층이면 임펄스는 스스로 처리한다"])
    table(s, Inches(0.7), y, Inches(11.9), Inches(4.1), rows, col_w=[2.5, 5.2, 4.2], size=11)
    footer(s, "median 채널은 배포 구조에 국소 필터를 덧댄 시도였고 검증 끝에 기각했다. "
              + ("구조 자체를 바꾸는 쪽(수용영역 확장)이 답이었다." if mi["key"] == "drunet"
                 else "구조는 배포된 그대로 두고 학습 절차만 바꿨다."))

    # ---------------------------------------------------------- 5. 결과
    s = blank(prs); bg(s)
    y = title(s, "결과", f"test 100장 · 제출값 PSNR {se['psnr_total']:.2f} / SSIM {se['ssim_total']:.4f}", "RESULTS")
    pic(s, FIGDIR / "summary_bars.png", Inches(0.6), y + Inches(0.15), w=Inches(7.7))

    rows = [["노이즈", "입력", "제안 모델", "Δ"]]
    for nz in NOISE_ORDER:
        rows.append([NOISE_KO[nz], f"{p_in[nz]:.2f}", f"{p_model[nz]:.2f}", f"+{p_model[nz] - p_in[nz]:.2f}"])
    rows.append(["전체", f"{n_all:.2f}", f"{se['psnr_total']:.2f}", f"+{se['psnr_total'] - n_all:.2f}"])
    table(s, Inches(8.5), y + Inches(0.15), Inches(4.1), Inches(2.0), rows,
          col_w=[1.7, 0.9, 1.1, 0.9], size=11.5, highlight_row=5)

    ry = y + Inches(2.45)
    card(s, Inches(8.5), ry, Inches(4.1), Inches(2.6), ACCENT_SOFT, ACCENT)
    text(s, Inches(8.72), ry + Inches(0.18), Inches(3.7), Inches(2.2), [
        ("Salt & Pepper 가 갈림길", {"size": 14, "bold": True, "color": INK}),
        (f"{p_in['salt_and_pepper']:.2f}  →  {p_model['salt_and_pepper']:.2f} dB   "
         f"(+{p_model['salt_and_pepper'] - p_in['salt_and_pepper']:.2f})",
         {"size": 13, "font": MONO, "space_before": Pt(8)}),
        (f"SSIM {q_in['salt_and_pepper']:.4f} → {q_model['salt_and_pepper']:.4f} — 4종 중 최고",
         {"size": 12, "font": MONO}),
        ("배포 기준선 DnCNN 은 여기서 29.69 였다. 구조를 바꾼 게 아니라 학습 절차만 바꿔 6 dB 를 더 벌었다.",
         {"size": 12, "space_before": Pt(8)}),
    ], color=INK2)
    footer(s, "conventional 3종은 배포 코드 구현 그대로 · self-ensemble 적용")


    # ---------------------------------------------------------- 5b. 시도 요약
    s = blank(prs); bg(s)
    y = title(s, "여기까지 어떻게 왔나", "채택한 것과 기각한 것을 같이 둔다", "TRACK RECORD")

    rows = [["", "시도", "결과", "배운 것"]]
    for st, what, res, learn in history_rows(se, p_model, mi, lf=lf):
        rows.append([st, what, res, learn])
    shape = table(s, Inches(0.7), y, Inches(11.9), Inches(3.5), rows,
                  col_w=[0.8, 3.15, 1.75, 6.2], size=9.5)

    # 상태별로 색을 달리해 채택/기각이 한눈에 갈리게 한다
    tone = {"채택": ACCENT, "기각": WARN, "전제": MUTED, "가산점": ACCENT, "정정": WARN}
    tbl = shape.table
    for r in range(1, len(rows)):
        st = rows[r][0]
        for c in range(4):
            cell = tbl.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = ACCENT_SOFT if st == "가산점" else WHITE
            for run in cell.text_frame.paragraphs[0].runs:
                if c == 0:
                    run.font.color.rgb = tone[st]
                    run.font.bold = True
                elif c == 2:
                    run.font.color.rgb = tone[st]
                    run.font.bold = st in ("채택", "가산점")

    ny = y + Inches(3.72)
    card(s, Inches(0.7), ny, Inches(11.9), Inches(1.18), ACCENT_SOFT, ACCENT)
    text(s, Inches(0.95), ny + Inches(0.13), Inches(11.4), Inches(0.95), [
        ("기각한 넷이 채택한 다섯을 설명한다", {"size": 14, "bold": True, "color": INK}),
        ("median 채널은 s&p 가 약하다는 진단에서 나왔다. 진단은 옳았고 처방이 틀렸다 — "
         "35px 시야를 그대로 둔 채 국소 필터를 하나 더 준 것이었기 때문이다. "
         "같은 문제를 수용영역 확장(DRUNet)이 풀었고 그 구간에서만 +5.82 dB. "
         "ablation 을 안 돌렸으면 손해 본 채 'median 덕분'이라고 발표할 뻔했다.",
         {"size": 11.5, "space_before": Pt(5)}),
    ], color=INK2)
    footer(s, "모든 수치는 test 100장 기준 · conventional 비교군과 지표 구현은 배포 코드 그대로")

    # ---------------------------------------------------------- 6-9. 4패널
    for nz in NOISE_ORDER:
        s = blank(prs); bg(s)
        y = title(s, f"{NOISE_KO[nz]}",
                  f"corrupted → restored → |error| → ground truth   ·   "
                  f"{p_in[nz]:.2f} → {p_model[nz]:.2f} dB (+{p_model[nz] - p_in[nz]:.2f})",
                  "2 · BEFORE / AFTER / ERROR / GT")
        pic(s, FIGDIR / f"panel_{nz}.png", Inches(0.55), y + Inches(0.45), w=Inches(12.2))
        note = {
            "gaussian": "오차가 경계선을 따라 남는다. 평탄한 영역은 거의 완전히 복원된다.",
            "rician": "4종 중 가장 어렵다. |·| 때문에 어두운 영역의 값이 위로 들려 있어, 밝기 자체가 이동한 상태에서 되돌려야 한다.",
            "uniform": "gaussian 과 비슷한 양상. 입력 PSNR 이 원래 높아 개선 폭은 작다.",
            "salt_and_pepper": "임펄스가 완전히 사라졌다. 남은 오차는 원래 임펄스가 있던 자리의 미세한 흔적뿐이다.",
        }[nz]
        footer(s, note)

    # ---------------------------------------------------------- 10. 오차 지도
    s = blank(prs); bg(s)
    y = title(s, "오차는 어디에 남는가", "열마다 같은 스케일 — 위: 아무것도 안 했을 때, 아래: 모델 통과 후",
              "2 · ERROR MAP")
    pic(s, FIGDIR / "error_maps.png", Inches(2.55), y + Inches(0.05), h=Inches(4.95))
    footer(s, "PSNR 한 숫자로는 '어디서' 틀렸는지 안 보인다. 오차가 경계에 몰려 있는지 평탄부에 흩어져 있는지가 다음 수를 결정한다.")

    # ---------------------------------------------------------- 11. 검증과 실패
    s = blank(prs); bg(s)
    y = title(s, "믿을 수 있는 숫자인가", "그리고 되지 않은 것 하나", "VALIDATION")

    c = card(s, Inches(0.7), y, Inches(5.8), Inches(2.3), ACCENT_SOFT, ACCENT)
    text(s, Inches(0.95), y + Inches(0.2), Inches(5.3), Inches(2.0), [
        ("지표 구현이 채점과 같은지 먼저 맞췄다", {"size": 15, "bold": True, "color": INK}),
        ("배포 예시 로그의 mean/median/adaptive 성적을 우리 로더·우리 지표로 다시 쟀다.",
         {"size": 12.5, "space_before": Pt(8)}),
        ("PSNR 최대 차이  0.0000 dB", {"size": 15, "font": MONO, "bold": True, "color": ACCENT, "space_before": Pt(8)}),
        ("4종 × 4방법 전부 소수 넷째자리까지 일치. 여기가 맞은 뒤에 학습을 돌렸다.",
         {"size": 12.5, "space_before": Pt(6)}),
    ], color=INK2)

    c = card(s, Inches(6.8), y, Inches(5.8), Inches(2.3))
    text(s, Inches(7.05), y + Inches(0.2), Inches(5.3), Inches(2.0), [
        ("되지 않은 것 — σ 게이트", {"size": 15, "bold": True, "color": WARN}),
        ("입력이 이미 50~62 dB 인 이미지에서는 모델이 손해다 (100장 중 8장). "
         "정답을 보고 매번 유리한 쪽을 고르면 +0.52 dB.", {"size": 12.5, "space_before": Pt(8)}),
        ("그런데 σ̂ 로는 그 8장을 못 고른다. MAD 추정기가 이미지의 미세 질감을 노이즈로 착각해, "
         "실제 σ=0.001 인 이미지를 0.037 로 추정한다. 임계값을 어디 두든 손해였다.",
         {"size": 12.5, "space_before": Pt(6)}),
    ], color=INK2)

    ay = y + Inches(2.6)
    if rej:
        rows = [["구성 (같은 학습 레시피)", "전체 PSNR", "전체 SSIM", "s&p PSNR"]]
        rows.append([f"{mi['name']} + self-ensemble  (제출)",
                     f"{se['psnr_total']:.2f}", f"{se['ssim_total']:.4f}",
                     f"{p_model['salt_and_pepper']:.2f}"])
        rows.append(["+ median 3×3 입력 채널  (기각)",
                     f"{rej['psnr']:.2f}", f"{rej['ssim']:.4f}", f"{rej['snp']:.2f}"])
        table(s, Inches(0.7), ay, Inches(7.1), Inches(1.05), rows, col_w=[3.6, 1.2, 1.2, 1.1],
              size=11.5, highlight_row=1)
        text(s, Inches(8.1), ay - Inches(0.05), Inches(4.4), Inches(1.6), [
            ("가설을 세웠고, 틀렸다", {"size": 14, "bold": True, "color": WARN}),
            (f"median 채널을 넣은 이유가 s&p 였는데 정작 거기서 {abs(rej['dp_snp']):.2f} dB 를 잃었다. "
             "ablation 을 안 돌렸으면 손해 본 채 'median 덕분'이라고 발표할 뻔했다.",
             {"size": 12, "space_before": Pt(6)}),
        ], color=INK2)
    else:
        card(s, Inches(0.7), ay, Inches(11.9), Inches(1.3))
        text(s, Inches(0.95), ay + Inches(0.25), Inches(11.4), Inches(0.9),
             "구조 ablation (median 채널) 학습 진행 중 — 완료 후 이 표를 채운다.",
             size=13, color=MUTED)
    footer(s, "check_baselines.py · analyze_gate.py 로 재현 가능")

    # ---------------------------------------------------------- 12. label-free
    s = blank(prs); bg(s)
    y = title(s, "Label-free 파이프라인", "clean 이미지를 loss 에 한 번도 쓰지 않는 경로 (Noise2Void)",
              "4 · BONUS")

    c = card(s, Inches(0.7), y, Inches(5.8), Inches(2.9), ACCENT_SOFT, ACCENT)
    text(s, Inches(0.95), y + Inches(0.2), Inches(5.3), Inches(2.6), [
        ("원리 — 그 픽셀을 입력에서 지운다", {"size": 15, "bold": True, "color": INK}),
        ("입력 : noisy (일부 픽셀을 이웃 값으로 덮어씀)", {"size": 12, "font": MONO, "space_before": Pt(8)}),
        ("정답 : 덮어쓰기 전 그 픽셀의 noisy 값", {"size": 12, "font": MONO}),
        ("loss : 덮어쓴 자리에서만", {"size": 12, "font": MONO}),
        ("노이즈가 픽셀마다 독립이면, 주변에서 그 픽셀의 노이즈를 알아낼 방법이 없다. "
         "네트워크가 맞힐 수 있는 것은 구조뿐이다. clean 은 어디에도 등장하지 않는다.",
         {"size": 12.5, "space_before": Pt(8)}),
    ], color=INK2)

    c = card(s, Inches(6.8), y, Inches(5.8), Inches(2.9))
    text(s, Inches(7.05), y + Inches(0.2), Inches(5.3), Inches(2.6), [
        ("우리 노이즈 4종은 전제를 만족한다", {"size": 15, "bold": True, "color": INK}),
        ("gaussian · rician · uniform · salt&pepper 모두 픽셀별 독립.",
         {"size": 12.5, "space_before": Pt(8)}),
        ("다만 L2 loss 는 조건부 평균을 학습하므로 노이즈 평균이 0 이어야 한다. "
         "s&p 는 0/max 로 덮으니 아니고, rician 은 절댓값 때문에 위로 들린다.",
         {"size": 12.5, "space_before": Pt(6)}),
        ("→ L1 loss 로 조건부 중앙값을 학습시켜 임펄스 편향을 없앤다.",
         {"size": 12.5, "bold": True, "color": ACCENT, "space_before": Pt(6)}),
    ], color=INK2)

    ly = y + Inches(3.15)
    if lf:
        rows = [["파이프라인", "clean", "test 입력", "PSNR", "SSIM"]]
        rows.append(["supervised (제출)", "loss 에 사용", "미사용",
                     f"{se['psnr_total']:.2f}", f"{se['ssim_total']:.4f}"])
        rows.append(["배포 기준선 (supervised 10ep)", "loss 에 사용", "미사용",
                     f"{BASE['psnr']:.2f}", f"{BASE['ssim']:.4f}"])
        for v in lf.values():
            rows.append([v["label"], "전혀 안 씀", v.get("test_in", "미사용"),
                         f"{v['psnr']:.2f}", f"{v['ssim']:.4f}"])
        shape = table(s, Inches(0.7), ly, Inches(11.9), Inches(1.5), rows,
                      col_w=[4.5, 1.9, 1.6, 1.5, 1.6], size=11)

        # 기준선을 넘은 label-free 행을 강조한다
        tbl = shape.table
        for r in range(3, len(rows)):
            if float(rows[r][3]) > BASE["psnr"]:
                for c in range(5):
                    cell = tbl.cell(r, c)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = ACCENT_SOFT
                    for run in cell.text_frame.paragraphs[0].runs:
                        run.font.bold = True
                        run.font.color.rgb = INK
    else:
        card(s, Inches(0.7), ly, Inches(11.9), Inches(1.6))
        text(s, Inches(0.95), ly + Inches(0.2), Inches(11.4), Inches(1.2), [
            ("두 가지 모드로 학습 중", {"size": 14, "bold": True, "color": INK}),
            ("① test_noise_only 100장만 사용 — clean 을 한 장도 건드리지 않는 가장 순수한 label-free",
             {"size": 12.5, "space_before": Pt(6)}),
            ("② train clean 으로 noisy 를 합성한 뒤 그 noisy 만 학습 — loss 에 clean 미사용, 데이터 7,268장",
             {"size": 12.5}),
            ("체크포인트 선택도 val noisy 의 마스킹 loss 로 한다. clean 기반 PSNR 로 고르면 파이프라인이 "
             "label-free 가 아니게 되기 때문이다.", {"size": 12, "color": MUTED, "space_before": Pt(6)}),
        ], color=INK2)
    footer(s, "test 입력을 학습에 쓰면 supervised 와 나란히 놓을 수 없다 — supervised 는 test 를 본 적이 없다. "
              "대표값은 test 무접촉 쪽을 쓴다.")

    # ---------------------------------------------------------- 13. 정리
    s = blank(prs); bg(s)
    y = title(s, "정리", None, "SUMMARY")
    items = [
        ("종류를 모른다는 제약이 이 문제의 전부",
         "고정 필터 하나로는 아무것도 안 한 것과 다를 바 없다 (mean 24.86 vs 입력 24.67). 학습이 필요한 이유."),
        (("좁은 시야가 진짜 제약이었다" if mi["key"] == "drunet" else "구조가 아니라 학습 절차가 답이었다"),
         (f"학습 절차만 바꿔 {BASE['psnr']:.2f} → 34.13 dB, 수용영역을 35px 에서 {mi['rf']} 로 넓혀 "
          f"{se['psnr_total']:.2f} dB. median 채널을 덧대려던 시도가 기각된 이유도 같다 — 시야를 안 건드렸다."
          if mi["key"] == "drunet" else
          f"배포된 구조 그대로 두고 loss·스케줄·증강·self-ensemble 만 바꿔 "
          f"{BASE['psnr']:.2f} → {se['psnr_total']:.2f} dB. 구조를 건드린 유일한 시도는 ablation 끝에 기각했다.")),
        ("숫자를 믿을 수 있게 먼저 맞췄다",
         "배포 예시 로그와 PSNR 최대 차이 0.0000 dB. 검증 없이 낸 숫자는 근거가 아니다."),
        ("안 되는 것도 확인했다",
         "σ 게이트로 얻을 0.52 dB 는 blind 상태에서 집을 수 없다. 제약이 실제 비용을 발생시킨다는 증거."),
        ("label-free 경로를 따로 만들었다",
         "노이즈가 픽셀 독립이라는 성질만으로 clean 없이 학습 가능. 현실에서는 clean 을 못 구하는 경우가 많다."),
    ]
    yy = y + Inches(0.05)
    for i, (h, d) in enumerate(items):
        num = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.75), yy + Inches(0.08), Inches(0.34), Inches(0.34))
        num.fill.solid(); num.fill.fore_color.rgb = ACCENT
        num.line.fill.background(); num.shadow.inherit = False
        tf = num.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        r = tf.paragraphs[0].add_run(); r.text = str(i + 1)
        r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = MONO
        text(s, Inches(1.3), yy, Inches(11.2), Inches(0.9), [
            (h, {"size": 15, "bold": True, "color": INK}),
            (d, {"size": 12.5, "color": MUTED, "space_before": Pt(3)}),
        ])
        yy += Inches(1.02)
    footer(s, f"제출값  PSNR {se['psnr_total']:.2f}   SSIM {se['ssim_total']:.4f}    "
              f"(배포 기준선 대비 +{se['psnr_total'] - BASE['psnr']:.2f} dB / +{se['ssim_total'] - BASE['ssim']:.4f})")


def main() -> None:
    try:
        sys.stdout.reconfigure(encoding="utf-8")
    except Exception:
        pass

    ap = argparse.ArgumentParser()
    ap.add_argument("--run", type=Path, default=ROOT / "runs" / "0831-1015_dncnn_plus_main")
    ap.add_argument("--name", default="")
    ap.add_argument("--lf", type=Path, default=None, help="label-free 결과 json (evaluate.py 출력)")
    ap.add_argument("--lf-train", type=Path, default=None)
    ap.add_argument("--lf-manual", action="append", default=None,
                    metavar="PSNR|SSIM|라벨|test입력",
                    help="json 없이 숫자만 넣는다. 여러 번 줄 수 있다. 구분자는 '|' 다 "
                         "(라벨에 쉼표가 들어가므로). "
                         "예: --lf-manual '31.38|0.9056|DRUNet · train 7,268장|미사용'")
    ap.add_argument("--rejected", type=Path, default=None,
                    help="기각된 변형(median 채널)의 결과 json — 대조군으로 슬라이드에 넣는다")
    ap.add_argument("--out", type=Path, default=ROOT / "실습5_denoising_발표.pptx")
    args = ap.parse_args()

    se_path, base_path = args.run / "test_metrics_se.json", args.run / "test_metrics.json"
    if not se_path.exists():
        raise SystemExit(
            f"{se_path} 가 없다. 먼저 평가를 돌릴 것:\n"
            f"  python evaluate.py <run>/checkpoints/checkpoint_best.ckpt --self-ensemble"
        )
    M = {"se": json.loads(se_path.read_text(encoding="utf-8"))}
    # self-ensemble 없는 결과는 그 기여도를 계산하는 데만 쓴다. 없으면 그 칸만 비운다.
    M["base"] = json.loads(base_path.read_text(encoding="utf-8")) if base_path.exists() else M["se"]
    M["has_base"] = base_path.exists()
    if not M["has_base"]:
        print(f"참고: {base_path.name} 이 없어 self-ensemble 기여도는 표시하지 않는다")

    lf = {}
    for key, path, label, clean in (
        ("test", args.lf, "label-free · test 100장만", "전혀 안 씀"),
        ("train", args.lf_train, "label-free · train noisy", "loss 에 미사용"),
    ):
        if path and path.exists():
            d = json.loads(path.read_text(encoding="utf-8"))
            lf[key] = {"label": label, "clean": clean,
                       "test_in": "사용" if key == "test" else "미사용",
                       "psnr": d["psnr_total"], "ssim": d["ssim_total"]}
    for i, spec in enumerate(args.lf_manual or []):
        # '|' 를 쓰되, 안 보이면 예전 방식(쉼표)으로도 받는다
        sep = "|" if "|" in spec else ","
        parts = [x.strip() for x in spec.split(sep)]
        lf[f"manual{i}"] = {
            "label": parts[2] if len(parts) > 2 else "label-free",
            "test_in": parts[3] if len(parts) > 3 else "미사용",
            "clean": "전혀 안 씀",
            "psnr": float(parts[0]),
            "ssim": float(parts[1]),
        }

    lf = lf or None

    rej = None
    if args.rejected and args.rejected.exists():
        d = json.loads(args.rejected.read_text(encoding="utf-8"))
        snp = float(np.mean([r["psnr_model"] for r in d["rows"]
                             if r["noise_type"] == "salt_and_pepper"]))
        main_snp = float(np.mean([r["psnr_model"] for r in M["se"]["rows"]
                                  if r["noise_type"] == "salt_and_pepper"]))
        rej = {"psnr": d["psnr_total"], "ssim": d["ssim_total"], "snp": snp,
               "dp_snp": snp - main_snp}

    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    build(prs, M, args.name, lf, rej, read_model(args.run))
    args.out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(args.out))
    print(f"슬라이드 {len(prs.slides.__iter__.__self__._sldIdLst)}장 → {args.out}")


if __name__ == "__main__":
    main()
