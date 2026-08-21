"""Pipeline + findings deck, rebuilt from the project's own data.

make_deck2.py lifted its images out of an older deck, so the figures could not follow
the analysis as it moved. This one composites every figure straight from the run
directories, which means the slides and FINDINGS.md cannot drift apart -- and it covers
what 4.7/4.8 added after that deck was made.

`person_3` is deliberately absent from every figure. Its reference set is selfies of an
identifiable real person, and this deck is built to be published; the concept still
appears in the numbers, just never as a face.

    python make_deck3.py --out 발표_파이프라인.pptx
"""

import argparse
import os
import sys

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.enum.text import PP_ALIGN

from make_deck2 import (
    BAND, CORAL, CORAL_SOFT, H, HEADER, INK, INK2, INK3, LINE, M, MANILLA, MONO,
    OCHRE, PAPER, SANS, W, WHITE, box, flow, inch, note, slide_base, table, text,
)

if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")

FIG = "deck3_img"

SKS = "v2_review/generated/{c}_{m}"
ZWX = "zwx10_review/zwx/{c}_{m}_zwx"
BASE = "baseline_review/generated/{c}_{m}"

FONTS = ["C:/Windows/Fonts/malgun.ttf", "C:/Windows/Fonts/arial.ttf",
         "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf"]


def font(size):
    for p in FONTS:
        if os.path.exists(p):
            return ImageFont.truetype(p, size)
    return ImageFont.load_default()


def square(path, side):
    with Image.open(path) as im:
        im = im.convert("RGB")
        s = min(im.size)
        l, t = (im.width - s) // 2, (im.height - s) // 2
        return im.crop((l, t, l + s, t + s)).resize((side, side), Image.LANCZOS)


def refs(concept, n=5):
    d = os.path.join("dataset", concept)
    fs = sorted(f for f in os.listdir(d) if f.lower().endswith((".jpg", ".jpeg", ".png")))
    return [os.path.join(d, f) for f in fs[:n]]


def grid(name, cells, cols, side=420, labels=None, quality=88):
    """cells: list of image paths (None -> blank). labels: parallel list of captions."""
    os.makedirs(FIG, exist_ok=True)
    out = os.path.join(FIG, name + ".jpg")
    cap = 34 if labels else 0
    rows = (len(cells) + cols - 1) // cols
    pad = 10
    Wp = pad + cols * (side + pad)
    Hp = pad + rows * (side + cap + pad)
    sheet = Image.new("RGB", (Wp, Hp), "white")
    d = ImageDraw.Draw(sheet)
    f = font(23)
    for i, p in enumerate(cells):
        r, c = divmod(i, cols)
        x = pad + c * (side + pad)
        y = pad + r * (side + cap + pad)
        if p and os.path.exists(p):
            sheet.paste(square(p, side), (x, y + cap))
        else:
            d.rectangle([x, y + cap, x + side, y + cap + side], fill="#eeeeee")
        if labels and i < len(labels) and labels[i]:
            d.text((x + 2, y + 4), labels[i], fill="#1F1E1D", font=f)
    sheet.save(out, "JPEG", quality=quality, optimize=True)
    return out


def strip(name, paths, side=300, quality=88):
    os.makedirs(FIG, exist_ok=True)
    out = os.path.join(FIG, name + ".jpg")
    sheet = Image.new("RGB", (side * len(paths), side), "white")
    for i, p in enumerate(paths):
        sheet.paste(square(p, side), (i * side, 0))
    sheet.save(out, "JPEG", quality=quality, optimize=True)
    return out


def build_figures():
    made = {}
    # the headline: same prompt, same everything, trigger token swapped
    made["trigger"] = grid("trigger", [
        SKS.format(c="actionfigure_2", m="lora") + "/0.png",
        ZWX.format(c="actionfigure_2", m="lora") + "/0.png",
        SKS.format(c="scene_waterfall", m="lora") + "/0.png",
        ZWX.format(c="scene_waterfall", m="lora") + "/0.png",
    ], cols=2, labels=["sks — 전투복 + 소총", "zwx", "sks — 폭포 앞 소총", "zwx"])

    made["refs_af"] = strip("refs_af", refs("actionfigure_2", 5))
    made["refs_tank"] = strip("refs_tank", refs("transport_tank", 4))

    # TI never learned the subject; LoRA did
    made["ti_fail"] = grid("ti_fail", [
        BASE.format(c="actionfigure_2", m="ti") + "/0.png",
        BASE.format(c="actionfigure_2", m="lora") + "/0.png",
    ], cols=2, labels=["TI  (DINO 0.14)", "LoRA  (DINO 0.55)"])

    # prior-preservation erases the instance
    made["prior"] = grid("prior", [
        SKS.format(c="luggage_backpack1", m="lora") + "/0.png",
        SKS.format(c="luggage_backpack1", m="lora_prior") + "/0.png",
    ], cols=2, labels=["prior 없음", "prior λ=1.0"])

    # background entanglement: desk-and-shallow-DOF look follows the trigger everywhere
    made["entangle"] = grid("entangle", [
        SKS.format(c="actionfigure_2", m="lora") + "/5.png",
        SKS.format(c="actionfigure_2", m="lora") + "/2.png",
    ], cols=2, labels=["'모래 해변'", "'설산 정상'"])

    # tank: the class noun beats the instance
    made["tank"] = grid("tank", [
        "dataset/transport_tank/" + sorted(os.listdir("dataset/transport_tank"))[0],
        SKS.format(c="transport_tank", m="lora") + "/0.png",
        ZWX.format(c="transport_tank", m="lora") + "/0.png",
    ], cols=3, labels=["레퍼런스 (용접 조형물)", "sks tank", "zwx tank"])
    return made


# --------------------------------------------------------------------------- slides

TOC = [
    ("01", "과제와 데이터"),
    ("02", "파이프라인"),
    ("03", "채점 방식"),
    ("04", "결과"),
    ("05", "지표는 믿을 만한가"),
    ("06", "개선 항목별 검증"),
    ("07", "핵심 발견 — 트리거 토큰"),
    ("08", "배운 점"),
]


def title_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = PAPER
    box(s, 0, 0, Pt(9), H, CORAL, None)
    text(s, M, inch(2.05), W - 2 * M, inch(0.3), "CUSTOMCONCEPT101 · SD3.5-MEDIUM",
         size=12.5, color=CORAL, bold=True, font=MONO)
    text(s, M, inch(2.5), W - 2 * M, inch(1.0),
         "이미지 생성 모델 개인화 —\n파이프라인과 지표 검증", size=40, bold=True)
    text(s, M, inch(4.12), inch(8.6), inch(0.9),
         "레퍼런스 몇 장으로 특정 대상을 학습시키고, 그 대상을 새로운 장면에 그린다.\n"
         "네 가지 학습 방식을 10개 컨셉에 적용하고, 무엇보다 그 결과를 재는 지표 자체를 검증했다.",
         size=15, color=INK2)
    box(s, M, inch(5.35), inch(4.4), Pt(1.5), LINE, None)
    text(s, M, inch(5.55), inch(8), inch(0.3),
         "TI · LoRA · DoRA · prior-preservation   |   CLIP T2I/I2I + DINO   |   10 컨셉 × 10 프롬프트",
         size=11.5, color=INK3, font=MONO)
    return s


def toc_slide(prs):
    s, y = slide_base(prs, "CONTENTS", "목차")
    half = (W - 2 * M - inch(0.6)) / 2
    for i, (num, name) in enumerate(TOC):
        col, row = divmod(i, 4)
        x = M + col * (Emu(int(half)) + inch(0.6))
        yy = y + inch(0.35) + row * inch(0.92)
        text(s, x, yy, inch(0.6), inch(0.3), num, size=15, color=CORAL, bold=True, font=MONO)
        text(s, x + inch(0.72), yy - inch(0.03), Emu(int(half)) - inch(0.8), inch(0.36),
             name, size=17, bold=True)
        box(s, x, yy + inch(0.46), Emu(int(half)) - inch(0.1), Pt(1), LINE, None)
    return s


def build(out, figs):
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H

    title_slide(prs)
    toc_slide(prs)

    # ---- 01 과제 -----------------------------------------------------------
    s, y = slide_base(prs, "01 · 과제와 데이터", "레퍼런스 몇 장 → 임의의 새 장면",
                      "개인화(personalization)란 특정 대상 하나를 모델에 가르쳐, 학습에 없던 상황에 그 대상을 그리게 하는 것이다.")
    s.shapes.add_picture(figs["refs_af"], M, y + inch(0.1), inch(7.5), inch(1.5))
    text(s, M, y + inch(1.72), inch(7.5), inch(0.3),
         "레퍼런스 6장 — 전부 책상 위 제품 사진", size=11, color=INK3)
    table(s, M + inch(8.1), y + inch(0.1), [
        ["항목", "값"],
        ["컨셉", "10"],
        ["컨셉당 레퍼런스", "4 ~ 15 장"],
        ["컨셉당 프롬프트", "10"],
        ["베이스 모델", "SD3.5-medium"],
        ["학습 스텝", "300"],
        ["총 생성 이미지", "1,000+"],
    ], [1.7, 1.5], row_h=0.34)
    note(s, y + inch(2.62),
         "어려운 이유는 두 목표가 서로 싸우기 때문이다. 대상을 레퍼런스와 똑같이 그릴수록 프롬프트가 요구하는\n"
         "새 장면에서 멀어지고, 장면을 충실히 그릴수록 대상이 흐려진다. 이 실험의 상당 부분은 그 균형을 재는 이야기다.")

    # ---- 02 파이프라인 -----------------------------------------------------
    s, y = slide_base(prs, "02 · 파이프라인", "데이터에서 결론까지 5단계",
                      "각 단계는 독립 스크립트이고, 출력 파일이 있으면 건너뛴다 — 중간에 죽어도 이어서 돌릴 수 있다.")
    flow(s, y + inch(0.18), [
        ("01", "데이터", "CustomConcept101\n컨셉별 레퍼런스 +\n프롬프트 10개"),
        ("02", "학습", "TI / LoRA / DoRA\n/ prior\n300 스텝, rank 16"),
        ("03", "생성", "프롬프트당 후보 4장\nCLIP 최고점 1장 선택\n512×512"),
        ("04", "평가", "CLIP T2I · I2I\n+ DINO\n프롬프트 단위 CSV"),
        ("05", "분석", "짝지은 부호검정\nHTML 뷰어\nFINDINGS.md"),
    ], h=1.62)
    table(s, M, y + inch(2.12), [
        ["스크립트", "역할"],
        ["run_all.py", "전 단계 오케스트레이션 · 재실행 안전"],
        ["train_textual_inversion.py / train_dreambooth_lora.py", "학습 (TI / LoRA · DoRA · prior)"],
        ["generate.py", "생성 + best-of-N 선택"],
        ["eval_breakdown.py", "CLIP·DINO 점수를 프롬프트 단위 CSV로"],
        ["make_viewer.py / zwx_table.py", "런 비교 뷰어 · 평가 표"],
    ], [5.6, 5.3], row_h=0.34)
    note(s, y + inch(4.34),
         "재실행 안전성이 실제로 값을 했다 — Colab 런타임이 끊겨 한 번 결과를 통째로 잃었고, 재학습으로 복구했다.")

    # ---- 03 채점 -----------------------------------------------------------
    s, y = slide_base(prs, "03 · 채점 방식", "지표 세 개, 서로 다른 것을 잰다")
    cols = [
        ("T2I", "텍스트 ↔ 생성물", "프롬프트를 얼마나 따랐나.\nCLIP 텍스트·이미지 코사인."),
        ("CLIP I2I", "레퍼런스 ↔ 생성물", "대상을 얼마나 닮았나.\nCLIP 이미지 임베딩 코사인."),
        ("DINO", "레퍼런스 ↔ 생성물", "정체성을 얼마나 지켰나.\n라벨 없이 이미지만으로 학습된\nViT-S/16 — DreamBooth 논문 지표."),
    ]
    bw = (W - 2 * M - inch(0.6)) / 3
    for i, (k, sub, body) in enumerate(cols):
        x = M + i * (Emu(int(bw)) + inch(0.3))
        box(s, x, y + inch(0.1), Emu(int(bw)), inch(1.95))
        box(s, x, y + inch(0.1), Emu(int(bw)), Pt(3), CORAL, None)
        text(s, x + inch(0.22), y + inch(0.34), Emu(int(bw)) - inch(0.44), inch(0.3),
             k, size=19, bold=True, font=MONO)
        text(s, x + inch(0.22), y + inch(0.68), Emu(int(bw)) - inch(0.44), inch(0.24),
             sub, size=11, color=CORAL, font=MONO)
        text(s, x + inch(0.22), y + inch(0.98), Emu(int(bw)) - inch(0.44), inch(0.9),
             body, size=12.5, color=INK2)
    note(s, y + inch(2.28),
         "DINO를 추가한 이유: CLIP은 이미지-텍스트 쌍으로 학습돼 '무엇이라고 말할지' 기준으로 임베딩이 조직된다.\n"
         "그래서 같은 클래스의 다른 개체를 잘 구분하지 못한다 — '어떤 배낭인가'보다 '배낭이다'를 본다.")
    table(s, M, y + inch(3.42), [
        ["", "TI", "LoRA", "방법 분리도"],
        ["CLIP I2I", "0.5654", "0.7105", "1.26 배"],
        ["DINO", "0.1604", "0.4628", ("2.89 배", "hi")],
    ], [3.0, 2.4, 2.4, 3.1], row_h=0.36)
    text(s, M, y + inch(4.68), W - 2 * M, inch(0.4),
         "'TI가 정체성을 학습하지 못했다'를 CLIP-I 0.57로는 애매하게, DINO 0.16으로는 명확하게 말할 수 있다.",
         size=13, color=INK2)

    # ---- 04 결과 -----------------------------------------------------------
    s, y = slide_base(prs, "04 · 결과", "방법별 성능 (10컨셉 평균)")
    table(s, M, y + inch(0.1), [
        ["방법", "T2I ↑", "CLIP I2I ↑", "DINO ↑"],
        ["Textual Inversion", "0.3040", "0.5876", ("0.2354", "hi")],
        ["LoRA", "0.3392", "0.7288", "0.5057"],
        ["DoRA", "0.3411", "0.7263", ("0.5192", "hi")],
        ["LoRA + prior", "0.3450", "0.6844", "0.4048"],
        ["DoRA + prior", "0.3437", "0.6830", "0.4009"],
    ], [4.4, 2.2, 2.2, 2.2], row_h=0.38)
    s.shapes.add_picture(figs["ti_fail"], M + inch(0.0), y + inch(2.62), inch(3.6), inch(1.94))
    text(s, M, y + inch(4.62), inch(4.2), inch(0.5),
         "TI는 배경만 프롬프트대로 그리고 대상 자체가 없다.", size=11.5, color=INK3)
    text(s, M + inch(4.1), y + inch(2.66), W - 2 * M - inch(4.3), inch(2.2),
         "- TI는 실패했다. DINO 0.24는 사실상 무관한 이미지를 뜻한다. 트리거 임베딩\n"
         "  하나(768차원)만 학습해서는 300스텝으로 개체를 못 담는다.\n"
         "- LoRA와 DoRA는 사실상 동률이다. CLIP I2I는 LoRA를, DINO는 DoRA를 앞세우는데\n"
         "  격차가 0.003 / 0.013으로 재학습 노이즈 안이다. 우열을 말할 수 없다.\n"
         "- prior-preservation은 정체성을 오히려 파괴했다 — 06장.",
         size=13, color=INK2)

    # ---- 05 지표 검증 ------------------------------------------------------
    s, y = slide_base(prs, "05 · 지표는 믿을 만한가", "두 지표는 프롬프트 단위로 서로 싸운다",
                      "컨셉 안 10개 프롬프트에 걸친 T2I–I2I 피어슨 상관 (baseline, 9컨셉 × 3방법)")
    table(s, M, y + inch(0.1), [
        ["방법", "평균 r", "음의 상관을 보인 컨셉"],
        ["TI", ("+0.187", "hi"), "3 / 9"],
        ["LoRA", "−0.415", "8 / 9"],
        ["DoRA", "−0.514", "8 / 9"],
    ], [3.0, 2.4, 3.4], row_h=0.38)
    note(s, y + inch(1.82),
         "정체성을 학습한 모델만 '레퍼런스처럼 보이기'와 '장면을 바꾸기' 사이의 경쟁에 직면한다.\n"
         "학습에 실패한 TI는 경쟁할 대상이 없어 두 지표가 무관하게 움직인다.\n"
         "→ 반상관의 존재 자체가 개인화 성공의 신호다. (부호검정 r(dora)<r(ti): 9/9, p=0.004)")
    text(s, M, y + inch(3.3), W - 2 * M, inch(0.32),
         "그런데 두 지표 모두 프롬프트별 체계적 교란을 갖는다", size=17, bold=True)
    table(s, M, y + inch(3.78), [
        ["교란", "증거", "함의"],
        ["CLIP I2I는 '구도 일치도'를 잰다",
         "person_3 레퍼런스가 전부 셀피 → 셀피 프롬프트 2개가 세 방법 모두에서 1·2위 (+31%)",
         "컨셉 간 I2I 비교 무효"],
        ["T2I는 프롬프트 길이에 끌린다",
         "가장 짧은 0번 프롬프트가 9컨셉 중 8컨셉에서 평균 이하",
         "프롬프트 간 T2I 비교 무효"],
    ], [3.5, 5.6, 2.6], row_h=0.42)
    text(s, M, y + inch(5.18), W - 2 * M, inch(0.3),
         "→ 유효한 비교는 '같은 프롬프트 안에서 방법끼리'뿐이다. 다행히 과제가 요구하는 비교가 정확히 그 형태다.",
         size=13, color=INK2)

    # ---- 06 ablation -------------------------------------------------------
    s, y = slide_base(prs, "06 · 개선 항목별 검증", "여섯 개를 바꿔보고, 넷은 효과가 없었다",
                      "검정 단위는 컨셉(컨셉당 1표). 짝지은 부호검정.")
    table(s, M, y + inch(0.1), [
        ["바꾼 것", "T2I", "정체성 (I2I / DINO)", "판정"],
        ["생성 해상도 256 → 512", "6/9 (p=0.51)", "5/9 (p=1.00)", "효과 없음"],
        ["rank 4 → 16", "기여 ≈ 0", "I2I +0.004~0.012", "미미"],
        ["학습 해상도 256 → 512", "부호 혼재", "부호 혼재, |Δ| ≤ 0.02", "측정 불가"],
        ["best-of-4 선택", ("+0.0089, 10/10 (p=0.002)", "hi"), "I2I −0.0035 (하락)", ("지표 게이밍", "hi")],
        ["prior-preservation λ=1.0", "+0.0067 (p=0.11, n.s.)", ("I2I −0.044 · DINO −0.101", "hi"), ("순손실", "hi")],
        ["prior λ=1.0 → 0.3", "유지", "손실의 29~66% 회복", "개선, 그러나 미달"],
    ], [3.9, 3.0, 3.2, 2.0], row_h=0.4)
    note(s, y + inch(3.06),
         "best-of-N의 선택 기준은 CLIP text-image 유사도 — 평가 지표 T2I와 완전히 같은 함수다.\n"
         "즉 추론 시점에 평가 지표를 직접 최적화하고 있고, 그 대가가 I2I 하락으로 관측된다.\n"
         "10/10 · p=0.002라는 가장 강한 결과가 사실은 가장 못 믿을 결과였다.", warn=True)
    s.shapes.add_picture(figs["prior"], M + inch(8.5), y + inch(4.4), inch(2.22), inch(1.2))
    text(s, M, y + inch(4.42), inch(8.2), inch(1.1),
         "prior-preservation은 클래스 붕괴를 막으려고 클래스 이미지 50장을 함께 학습시킨다.\n"
         "그런데 λ가 기본 1.0이라 그래디언트 예산의 절반이 거기 쓰인다 — 300스텝 예산에서는 과하다.\n"
         "DINO가 CLIP-I보다 2.3배 크게 떨어진 것이 결정적이다: 구도 아티팩트였다면 DINO는 버텼어야 한다.",
         size=12.5, color=INK2)

    # ---- 07 트리거 (핵심) ---------------------------------------------------
    s, y = slide_base(prs, "07 · 핵심 발견", "가장 큰 결함은 하이퍼파라미터가 아니라 단어 하나였다")
    text(s, M, y + inch(0.06), inch(6.3), inch(1.5),
         "DreamBooth 관행을 따라 트리거 토큰을 `sks`로 잡았다.\n"
         "그런데 sks는 의미 없는 문자열이 아니라 SKS 소총이고,\n"
         "SD3.5의 텍스트 인코더는 그렇게 읽는다.\n\n"
         "학습 문장이 \"a photo of sks {클래스}\" 하나뿐이라,\n"
         "소총 prior가 매 스텝 트리거에 함께 묶인다.",
         size=14, color=INK2)
    box(s, M, y + inch(1.82), inch(6.3), inch(1.34), MANILLA, None)
    text(s, M + inch(0.24), y + inch(1.98), inch(5.9), inch(1.1),
         "10컨셉 0번 프롬프트를 눈으로 세어보면\n"
         "  · 명확한 오염 3개 — 소총·전투복이 프레임에 등장\n"
         "  · 약한 오염 1개 — 올리브색 밀리터리 재킷\n"
         "  · 나머지 6개는 육안상 깨끗",
         size=13, color=OCHRE)
    s.shapes.add_picture(figs["trigger"], M + inch(6.9), y + inch(0.02), inch(3.85), inch(3.9))
    text(s, M + inch(6.9), y + inch(3.96), inch(4.4), inch(0.3),
         "같은 프롬프트 · 같은 설정 · 트리거만 교체", size=11, color=INK3)
    note(s, y + inch(3.34),
         "사용자가 요청한 적 없는 소총이 이미지에 들어온다. 이건 성능 저하가 아니라 결함이다.")

    # ---- 07b zwx 결과 ------------------------------------------------------
    s, y = slide_base(prs, "07 · 핵심 발견", "고쳤더니 — 지표가 오히려 벌을 줬다",
                      "10컨셉 × LoRA/DoRA 재학습. rank·해상도·best-of-N 동일, 트리거 토큰만 교체.")
    table(s, M, y + inch(0.1), [
        ["그룹", "ΔT2I", "ΔCLIP I2I", "ΔDINO (LoRA)", "ΔDINO (DoRA)"],
        ["오염 4컨셉", "+0.0059 (3/4)", "+0.0249 (3/4)", "−0.0014 (1/4)", "+0.0013 (3/4)"],
        ["깨끗한 5컨셉", "+0.0050 (3/5)", "−0.0065 (2/5)", "−0.0207 (1/5)", "−0.0323 (2/5)"],
        ["전체 10컨셉", "+0.0050 (7/10)", "+0.0051 (5/10)", ("−0.0155 (2/10)", "hi"), "−0.0282 (5/10)"],
    ], [2.6, 2.3, 2.3, 2.4, 2.4], row_h=0.4)
    text(s, M, y + inch(1.9), W - 2 * M, inch(0.32),
         "컨셉 단위에서는 어떤 지표도 유의하지 않다 — 어느 방향으로도. 그러나 방향이 예측과 반대다.",
         size=15, bold=True)
    text(s, M, y + inch(2.34), inch(6.4), inch(1.5),
         "- 3컨셉 단계에서 '깨끗한 컨셉엔 무해할 것'이라 예측했다. 틀렸다 —\n"
         "  LoRA 기준 10컨셉 중 8개에서 DINO가 하락했다.\n"
         "- 이득은 사실상 한 컨셉에 몰려 있다 (actionfigure_2 +0.093).\n"
         "  나머지 하락 8개는 전부 −0.05 이내로 재학습 노이즈와 구분되지 않는다.\n"
         "- 즉 방향만 확인됐고 크기는 측정되지 않았다.",
         size=13, color=INK2)
    box(s, M + inch(6.9), y + inch(2.3), inch(3.85), inch(1.72), CORAL_SOFT, None)
    text(s, M + inch(7.14), y + inch(2.46), inch(3.4), inch(1.4),
         "검정 단위를 조심해야 한다\n\n"
         "한 컨셉의 프롬프트 10개는 같은 가중치에서\n나오므로 독립이 아니다.\n\n"
         "DoRA ΔDINO — 프롬프트 100개: p=0.002\n"
         "                    컨셉 10개: p=1.000",
         size=12, color=INK2)
    note(s, y + inch(4.22),
         "그래도 zwx를 써야 한다. 단, 지표를 근거로 들 수는 없다 — 근거는 '사용자가 요청하지 않은 소총이 사라진다'는 사실뿐이다.")

    # ---- 07c 지표 실패 정리 -------------------------------------------------
    s, y = slide_base(prs, "07 · 핵심 발견", "지표가 실패하는 세 가지 방식",
                      "이 실험의 결론은 '어떤 방법이 좋은가'보다 '무엇으로 재는가'에 있다.")
    rows = [
        ("게이밍당한다", "best-of-4",
         "선택기와 평가기가 같은 함수라 T2I가 10/10 올랐다.\n품질 개선이 아니라 지표 최적화였다."),
        ("보지 못한다", "학습 해상도 512",
         "CLIP도 DINO도 입력을 224로 줄인다.\n512 학습이 만든 디테일은 지표에 물리적으로 닿지 않는다."),
        ("반대로 채점한다", "트리거 토큰",
         "명백한 결함을 고쳤는데 DINO가 8/10 컨셉에서 하락했다.\n세 지표를 다 써도 '폭포 앞 소총'을 잡지 못한다."),
    ]
    bw = (W - 2 * M - inch(0.6)) / 3
    for i, (k, case, body) in enumerate(rows):
        x = M + i * (Emu(int(bw)) + inch(0.3))
        box(s, x, y + inch(0.2), Emu(int(bw)), inch(2.5))
        box(s, x, y + inch(0.2), Emu(int(bw)), Pt(3), CORAL, None)
        text(s, x + inch(0.24), y + inch(0.5), Emu(int(bw)) - inch(0.48), inch(0.34),
             k, size=18, bold=True)
        text(s, x + inch(0.24), y + inch(0.92), Emu(int(bw)) - inch(0.48), inch(0.24),
             case, size=11.5, color=CORAL, font=MONO)
        text(s, x + inch(0.24), y + inch(1.28), Emu(int(bw)) - inch(0.48), inch(1.2),
             body, size=12.5, color=INK2)
    note(s, y + inch(2.98),
         "재학습 노이즈 바닥도 함께 재야 했다. 유실된 런을 복구하며 같은 설정을 두 번 학습한 셈이 됐고,\n"
         "그 차이가 평균 |ΔI2I| 0.019 · 최대 0.041이었다 — 이 문서가 보고한 여러 효과와 같은 크기다.\n"
         "→ 컨셉 3개짜리 학습측 실험에서 0.02 이하의 차이는 보고할 근거가 없다.")

    # ---- 08 남은 실패 ------------------------------------------------------
    s, y = slide_base(prs, "08 · 배운 점", "고쳐지지 않은 것들")
    s.shapes.add_picture(figs["tank"], M, y + inch(0.12), inch(5.4), inch(1.94))
    text(s, M, y + inch(2.12), inch(5.6), inch(0.6),
         "클래스 명사 충돌 — 레퍼런스는 자전거 체인으로 만든 손바닥 크기 조형물인데\n"
         "class_prompt가 \"tank\"다. T2I는 진짜 군용 탱크를 그릴수록 오르고,\n"
         "DINO는 조형물을 닮을수록 오른다. 동시 최대화가 불가능하다.",
         size=12, color=INK2)
    s.shapes.add_picture(figs["entangle"], M + inch(5.9), y + inch(0.12), inch(4.0), inch(2.16))
    text(s, M + inch(5.9), y + inch(2.34), inch(4.9), inch(0.6),
         "배경 entanglement — '모래 해변'·'설산 정상'을 요구해도 레퍼런스의\n"
         "'어두운 책상 위 얕은 심도' 룩을 끌고 다닌다. 인스턴스 프롬프트가\n"
         "하나뿐이라 배경·조명·앵글이 전부 트리거에 흡수된다.",
         size=12, color=INK2)
    table(s, M, y + inch(3.36), [
        ["앞으로", "왜"],
        ["선택기와 평가기 분리 (DINO로 best-of-N)", "지표 게이밍 제거. 추론 전용이라 재학습 불필요"],
        ["이미지별 캡션 학습", "배경을 텍스트가 설명하면 트리거가 외울 필요가 없다"],
        ["시드 3~5회 반복", "8/10 DINO 하락이 실재인지 노이즈인지 가릴 유일한 방법"],
        ["TI·prior 변형도 zwx로 통일", "현재 결과표가 두 트리거에 걸쳐 있다"],
    ], [5.4, 5.5], row_h=0.36)

    # ---- 09 산출물 ---------------------------------------------------------
    s, y = slide_base(prs, "08 · 배운 점", "한 줄 요약과 산출물")
    box(s, M, y + inch(0.1), W - 2 * M, inch(1.16), CORAL_SOFT, None)
    text(s, M + inch(0.36), y + inch(0.3), W - 2 * M - inch(0.72), inch(0.9),
         "지표는 결론을 검증하는 도구가 아니라, 그 자체가 검증 대상이었다.\n"
         "가장 강한 통계(10/10, p=0.002)가 가장 못 믿을 결과였고, 가장 명백한 결함은 어떤 지표에도 잡히지 않았다.",
         size=16, bold=True, color=INK)
    table(s, M, y + inch(1.56), [
        ["산출물", "내용"],
        ["FINDINGS.md", "전체 분석 — 지표 검증 · ablation 8건 · 정성 관찰 · 한계"],
        ["view/viewer_*.html", "컨셉별 뷰어. 4개 런 × 최대 11개 방법을 프롬프트 단위로 나란히"],
        ["zwx10_review/zwx_table.md", "트리거 토큰 평가 표 (컨셉 단위 검정 포함)"],
        ["*_review/breakdown_*.csv", "프롬프트 단위 원점수 — T2I · I2I · DINO"],
        ["run_all.py 외 스크립트 8종", "재실행 안전한 학습 · 생성 · 평가 · 시각화 파이프라인"],
    ], [4.0, 6.9], row_h=0.38)
    text(s, M, y + inch(4.06), W - 2 * M, inch(0.3),
         "person_3은 실존 인물의 셀피가 레퍼런스이므로 이 발표 자료의 모든 그림에서 제외했다 — 수치로만 등장한다.",
         size=12, color=INK3)

    prs.save(out)
    return out, len(prs.slides._sldIdLst)


if __name__ == "__main__":
    ap = argparse.ArgumentParser()
    ap.add_argument("--out", default="발표_파이프라인.pptx")
    args = ap.parse_args()

    figs = build_figures()
    print(f"figures: {len(figs)} -> {FIG}/")
    out, n = build(args.out, figs)
    print(f"[wrote] {out}  ({os.path.getsize(out) / 1e6:.1f} MB, {n} slides)")
