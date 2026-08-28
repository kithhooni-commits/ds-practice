"""얼굴 복원 발표 슬라이드(.pptx) 생성.

웹 자료 두 장(468점에서 얼굴까지 / 468점의 한계)을 같은 팔레트·구조로 옮긴다.
수치는 전부 실측값이다 — 임의로 바꾸지 말 것.

    python docs/make_face_ppt.py
"""
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Pt

OUT = Path(__file__).parent / "얼굴복원_발표.pptx"

# ── 팔레트 (웹 자료와 동일) ──────────────────────────────────────────
PAPER  = RGBColor(0xF4, 0xF7, 0xF7)
CARD   = RGBColor(0xFF, 0xFF, 0xFF)
INK    = RGBColor(0x12, 0x20, 0x1F)
MUTED  = RGBColor(0x5C, 0x6B, 0x69)
ACCENT = RGBColor(0x0E, 0x8F, 0x86)
SIGNAL = RGBColor(0xC2, 0x4A, 0x2F)
RULE   = RGBColor(0xD8, 0xE0, 0xDE)
WASH   = RGBColor(0xE9, 0xF0, 0xEF)

BODY = "맑은 고딕"
MONO = "Consolas"

# 16:9
W, H = Emu(12192000), Emu(6858000)
M = Emu(640080)                      # 좌우 여백


def inch(v):
    return Emu(int(v * 914400))


def box(slide, x, y, w, h, fill=None, line=None, line_w=1.0):
    """색 채운 사각형. fill/line 은 RGBColor 또는 None."""
    from pptx.enum.shapes import MSO_SHAPE
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sh.shadow.inherit = False
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(line_w)
    sh.text_frame.word_wrap = True
    return sh


def text(slide, x, y, w, h, runs, *, size=14, color=INK, bold=False,
         font=BODY, align=PP_ALIGN.LEFT, spacing=1.15, anchor=MSO_ANCHOR.TOP):
    """runs: 문자열 또는 [(문자열, {size,color,bold,font}), ...] 목록의 목록(문단)."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

    # 입력 형태를 문단 목록으로 정규화한다.
    #   "글"                     -> 문단 1개, 런 1개
    #   ["글", "글"]             -> 문단 2개
    #   [("글", {...}), ...]     -> 문단 1개, 런 여러 개   ← 튜플로 시작하면 런 목록
    #   [[("글", {...})], ...]   -> 문단 여러 개, 각각 런 목록
    if isinstance(runs, str):
        paras = [runs]
    elif runs and isinstance(runs[0], tuple):
        paras = [runs]                      # 런 목록 하나 = 문단 하나
    else:
        paras = runs
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        chunks = [(para, {})] if isinstance(para, str) else para
        for s, opt in chunks:
            r = p.add_run()
            r.text = s
            f = r.font
            f.name = opt.get("font", font)
            f.size = Pt(opt.get("size", size))
            f.bold = opt.get("bold", bold)
            f.color.rgb = opt.get("color", color)
    return tb


def slide_base(prs, eyebrow, title, sub=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = PAPER

    text(s, M, inch(0.42), W - 2 * M, inch(0.25),
         [(eyebrow, {"font": MONO, "size": 10.5, "color": ACCENT, "bold": True})])
    text(s, M, inch(0.70), W - 2 * M, inch(0.62),
         [(title, {"size": 30, "bold": True, "color": INK})])
    if sub:
        text(s, M, inch(1.34), inch(9.2), inch(0.5),
             [sub], size=13, color=MUTED, spacing=1.25)
    # 강조 밑줄
    box(s, M, inch(1.22), inch(0.42), Emu(28575), fill=ACCENT)
    return s


def spec_row(s, items, y, h=inch(0.92)):
    """균등 분할 카드 행. items: [(라벨, 값, 보조), ...]"""
    n = len(items)
    gap = inch(0.10)
    total = W - 2 * M
    cw = Emu(int((total - gap * (n - 1)) / n))
    for i, (k, v, note) in enumerate(items):
        x = Emu(int(M + i * (cw + gap)))
        box(s, x, y, cw, h, fill=CARD, line=RULE)
        text(s, Emu(x + inch(0.16)), Emu(y + inch(0.13)), Emu(cw - inch(0.3)), inch(0.2),
             [(k, {"font": MONO, "size": 9, "color": MUTED, "bold": True})])
        text(s, Emu(x + inch(0.16)), Emu(y + inch(0.33)), Emu(cw - inch(0.3)), inch(0.3),
             [(v, {"font": MONO, "size": 19, "color": INK, "bold": True})])
        if note:
            text(s, Emu(x + inch(0.16)), Emu(y + inch(0.66)), Emu(cw - inch(0.3)), inch(0.22),
                 [(note, {"size": 10, "color": MUTED})])


def step_row(s, no, title, desc, shape, y, h=inch(0.80)):
    """파이프라인 한 단계: 번호 · 제목/설명 · 데이터 형태"""
    box(s, M, y, W - 2 * M, Emu(9525), fill=RULE)       # 상단 괘선
    yy = Emu(y + inch(0.12))
    text(s, M, yy, inch(0.4), inch(0.25),
         [(no, {"font": MONO, "size": 11, "color": ACCENT, "bold": True})])
    tx = Emu(M + inch(0.45))
    text(s, tx, yy, inch(5.4), inch(0.24), [(title, {"size": 13, "bold": True})])
    text(s, tx, Emu(yy + inch(0.26)), inch(5.4), inch(0.46),
         [desc], size=10.5, color=MUTED, spacing=1.2)
    # 데이터 형태 칩
    sx = Emu(M + inch(6.15))
    sw = inch(3.55)
    box(s, sx, yy, sw, inch(0.56), fill=WASH, line=None)
    box(s, sx, yy, Emu(19050), inch(0.56), fill=ACCENT)
    text(s, Emu(sx + inch(0.14)), Emu(yy + inch(0.09)), Emu(sw - inch(0.24)), inch(0.4),
         shape, size=10.5, font=MONO, color=INK, spacing=1.2)


def card(s, x, y, w, h, tag, tag_color, title, body):
    box(s, x, y, w, h, fill=CARD, line=RULE)
    text(s, Emu(x + inch(0.18)), Emu(y + inch(0.14)), Emu(w - inch(0.36)), inch(0.2),
         [(tag, {"font": MONO, "size": 8.5, "color": tag_color, "bold": True})])
    text(s, Emu(x + inch(0.18)), Emu(y + inch(0.36)), Emu(w - inch(0.36)), inch(0.26),
         [(title, {"size": 12.5, "bold": True})])
    text(s, Emu(x + inch(0.18)), Emu(y + inch(0.66)), Emu(w - inch(0.36)), Emu(h - inch(0.8)),
         [body], size=10.5, color=MUTED, spacing=1.22)


def evidence(s, x, y, w, chunks, h=inch(0.52)):
    """증거 줄 — 왼쪽에 강조 막대"""
    box(s, x, y, w, h, fill=WASH, line=None)
    box(s, x, y, Emu(19050), h, fill=ACCENT)
    text(s, Emu(x + inch(0.14)), Emu(y + inch(0.09)), Emu(w - inch(0.26)), Emu(h - inch(0.16)),
         chunks, size=10, font=MONO, color=INK, spacing=1.25)


# ══════════════════════════════════════════════════════════════════
def build():
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H

    # ── 1. 표지 ───────────────────────────────────────────────────
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = PAPER
    text(s, M, inch(2.05), W - 2 * M, inch(0.3),
         [("COMPUTER VISION 실습 · 얼굴 복원", {"font": MONO, "size": 12, "color": ACCENT, "bold": True})])
    text(s, M, inch(2.45), W - 2 * M, inch(1.1),
         [("468점에서 얼굴까지", {"size": 54, "bold": True})])
    box(s, M, inch(3.62), inch(0.62), Emu(38100), fill=ACCENT)
    text(s, M, inch(3.86), inch(8.4), inch(0.8),
         ["웹캠 한 프레임이 3D 얼굴이 되기까지, 그리고 그 과정에서 배운 것"],
         size=17, color=MUTED, spacing=1.3)
    text(s, M, inch(5.85), inch(9.0), inch(0.3),
         [[("MediaPipe Face Mesh · Three.js · 브라우저 단독", {"font": MONO, "size": 10.5, "color": MUTED})]])

    # ── 2. 쓴 모델 ────────────────────────────────────────────────
    s = slide_base(prs, "MODEL", "쓴 모델",
                   "학습은 하지 않았다 — 사전학습 모델의 출력 그 자체가 메시이고, 나머지는 전부 기하 처리다.")
    spec_row(s, [
        ("모델",  "Face Mesh", "Google · 사전학습"),
        ("구성",  "검출 + 회귀", "얼굴 검출기 + 랜드마크망"),
        ("출력",  "468 × 3", "(x, y, z)"),
        ("토폴로지", "고정", "번호 = 같은 지점"),
        ("실행",  "브라우저", "서버 · GPU 없음"),
        ("우리 학습", "없음", "추론만"),
    ], inch(2.05))
    text(s, M, inch(3.20), inch(10.9), inch(0.9),
         ["좌표는 정규화되어 나온다 — x는 이미지 너비로, y는 높이로. z는 머리 중심 기준 상대 깊이이고 "
          "x와 대략 같은 축척이다. 옵션 하나(refineLandmarks)를 켜면 홍채 10점이 붙어 478점이 된다."],
         size=12.5, color=MUTED, spacing=1.35)

    card(s, M, inch(4.05), inch(5.35), inch(1.55), "핵심 1", ACCENT,
         "번호가 곧 대응점이다",
         "토폴로지가 고정이라 어느 각도·어느 사람이든 17번은 항상 같은 지점. "
         "다시점 복원에서 가장 어려운 대응점 찾기가 통째로 사라진다.")
    card(s, Emu(M + inch(5.55)), inch(4.05), inch(5.35), inch(1.55), "핵심 2", ACCENT,
         "정합이 공짜다",
         "사진과 좌표가 같은 프레임에서 나온다. 화면 좌표를 그대로 UV로 쓰면 "
         "텍스처가 정확히 맞는다 — 얼굴 정렬 단계가 아예 없다.")

    # ── 3. 파이프라인 ─────────────────────────────────────────────
    s = slide_base(prs, "PIPELINE", "파이프라인 6단계",
                   "오른쪽이 각 단계의 데이터 형태다. 숫자가 어떻게 변해 가는지 따라가면 된다.")
    y0 = inch(1.90)
    steps = [
        ("01", "촬영 · 평균", "12프레임의 랜드마크를 평균낸다. 한 프레임만 쓰면 추정 지터가 그대로 표면 굴곡이 된다.",
         [[("12 × (468 × 3)", {})], [("→ 468 × 3  (평균)", {"color": MUTED})]]),
        ("02", "삼각형 복원", "토폴로지는 간선 쌍 목록이다. 3개씩 묶으면 한 삼각형 — 전부 닫히고 어긋난 그룹 0개.",
         [[("2,556 간선쌍", {})], [("→ 852 삼각형", {"color": MUTED})]]),
        ("03", "구멍 메우기", "테셀레이션은 눈꺼풀·입술의 둘레만 덮는다. 경계 간선을 이어 루프를 찾아 부채꼴로 메운다.",
         [[("경계 루프 4개", {})], [("얼굴선 36 · 입 20 · 눈 16·16", {"color": MUTED, "size": 9})],
          [("→ 904 삼각형", {"color": MUTED})]]),
        ("04", "정규화", "x에 종횡비를 곱해 축별 정규화를 되돌리고, 무게중심·축척을 맞춰 카메라 거리에 무관하게 만든다.",
         [[("x ×= W/H", {})], [("480×360 에서 0.75배 눌림 보정", {"color": MUTED, "size": 9})]]),
        ("05", "텍스처 · UV", "랜드마크의 화면 좌표가 곧 UV다. 같은 프레임에서 나왔으므로 정합 과정이 없다.",
         [[("512 × 512 JPEG (~50KB)", {})], [("+ 468 × 2 UV", {"color": MUTED})]]),
        ("06", "변형 · 전송", "피격은 감쇠 진동으로 눌렀다 되돌린다. 네트워크로는 좌표와 텍스처만 보낸다.",
         [[("전송 ≈ 60KB / 1회", {})], [("촬영할 때 한 번뿐", {"color": MUTED, "size": 9})]]),
    ]
    for i, (no, t, d, sh) in enumerate(steps):
        step_row(s, no, t, d, sh, Emu(y0 + i * inch(0.78)))
    box(s, M, Emu(y0 + 6 * inch(0.78)), W - 2 * M, Emu(9525), fill=RULE)

    # ── 4. 무엇이 아닌가 ──────────────────────────────────────────
    s = slide_base(prs, "SCOPE", "무엇이 아닌가",
                   "질문이 나오기 쉬운 지점. 여기서 한계의 이유가 그대로 설명된다.")
    card(s, M, inch(2.00), inch(5.35), inch(1.70), "아닌 것", SIGNAL,
         "3DMM 피팅이 아니다",
         "통계 얼굴 모델(FLAME·BFM)의 계수를 최적화하는 방식이 아니다. "
         "그래서 빠르지만 두개골·귀·목이 없다 — 모델이 담고 있지 않기 때문이다.")
    card(s, Emu(M + inch(5.55)), inch(2.00), inch(5.35), inch(1.70), "아닌 것", SIGNAL,
         "포토그래메트리가 아니다",
         "여러 장에서 깊이를 삼각측량하지 않는다. z는 모델이 추정한 값이라 "
         "옆 얼굴 깊이는 관측이 아니라 예측이다.")
    text(s, M, inch(4.05), inch(10.9), inch(0.4),
         [[("그래서 뒤통수는 지어낼 수 없다.", {"size": 17, "bold": True, "color": INK})]])
    text(s, M, inch(4.50), inch(10.9), inch(0.9),
         ["얼굴 테두리에서 두개골을 만들어 붙여 봤지만, 사람마다 테두리 모양과 깊이가 크게 달라 "
          "상수로 맞출 수 있는 문제가 아니었다. 비율은 맞췄으나(깊이/폭 0.70 → 1.04) 실제 얼굴에서 "
          "형태가 무너졌다. 지금은 얼굴을 등록하면 아바타 머리를 숨기고 얼굴만 남긴다."],
         size=12.5, color=MUTED, spacing=1.35)

    # ── 5. 배운 것 (1) ────────────────────────────────────────────
    s = slide_base(prs, "LESSONS 1/2", "배운 것 — 검증",
                   "눈으로 보고 짐작한 것은 세 번 연속 틀렸다. 회색 줄이 각 항목의 근거다.")
    y = inch(1.95)
    for no, t, d, ev in [
        ("01", "좌표를 세는 검사로는 못 잡는 버그가 있다",
         "뒤통수 텍스처를 옆 사진에서 가져오게 만들고 UV 검사를 전부 통과시켰다. "
         "그런데 화면에는 천장과 벽이 뒤통수에 발려 있었다. UV가 옆면 칸 안인 것은 맞았다 — "
         "그 칸의 그 자리가 머리가 아니었을 뿐이다.",
         [[("UV 검사 4/4 통과 · 실제 뒤통수 픽셀의 ", {}), ("79.7%", {"color": SIGNAL, "bold": True}),
           ("가 사진 배경", {})],
          [("→ 머리 바깥을 미리 머리카락색으로 덮음 · 오염 ", {"color": MUTED}),
           ("0%", {"color": ACCENT, "bold": True})]]),
        ("02", "좌표계 불일치는 예외를 던지지 않는다",
         "랜드마크는 x를 너비로, y를 높이로 정규화해 비정방형 프레임에서 얼굴이 눌린다. "
         "텍스처는 잘라낸 사진인데 랜드마크는 전체 프레임 기준이라 머리 타원이 1/4로 잡힌다. "
         "어느 쪽도 에러가 나지 않고 그럴듯한 그림이 나온다.",
         [[("480×360 에서 가로 ", {}), ("0.75배", {"color": SIGNAL, "bold": True}), (" 눌림", {})],
          [("머리 바깥 덮기가 얼굴을 삼킴 → crop 통과 후 얼굴 ", {"color": MUTED}),
           ("63%", {"color": ACCENT, "bold": True}), (" 생존", {"color": MUTED})]]),
    ]:
        box(s, M, y, W - 2 * M, Emu(9525), fill=RULE)
        text(s, M, Emu(y + inch(0.14)), inch(0.4), inch(0.25),
             [(no, {"font": MONO, "size": 11, "color": ACCENT, "bold": True})])
        tx = Emu(M + inch(0.45))
        text(s, tx, Emu(y + inch(0.12)), inch(10.3), inch(0.26), [(t, {"size": 14, "bold": True})])
        text(s, tx, Emu(y + inch(0.44)), inch(10.3), inch(0.62), [d], size=11, color=MUTED, spacing=1.28)
        evidence(s, tx, Emu(y + inch(1.12)), inch(10.2), ev)
        y = Emu(y + inch(1.86))

    # ── 6. 배운 것 (2) ────────────────────────────────────────────
    s = slide_base(prs, "LESSONS 2/2", "배운 것 — 데이터와 성능", None)
    y = inch(1.70)
    for no, t, d, ev in [
        ("03", "합성 데이터는 만점을 준다",
         "손으로 만든 궤적·토폴로지로 검사하면 그 데이터에 맞춰 만든 임계값이 당연히 맞는다. "
         "회귀 검사로는 쓸모가 있지만 정확도의 근거는 못 된다. 얼굴에서도 같았다 — "
         "합성 토폴로지에는 배경이 없어 배경 버그가 잡힐 수 없었다.",
         [[("합성 11케이스 F1 ", {}), ("1.000", {"color": ACCENT, "bold": True}),
           ("   ·   라벨 붙인 실영상 90초 F1 ", {}), ("0.386", {"color": SIGNAL, "bold": True})]]),
        ("04", "성능은 단독으로 재면 안 된다",
         "더 나은 복원을 위해 생성 모델(TripoSR)을 붙여 봤다. 단독 실행은 충분히 빨랐는데 "
         "게임 서버 안에서는 끝나지 않았다. GPU 경합인 줄 알았지만 아니었다 — 게임 루프가 "
         "초당 수십 번 브로드캐스트하며 GIL을 쥐고 있어 추론 스레드가 계속 굶었다.",
         [[("단독 ", {}), ("16.4s", {"color": ACCENT, "bold": True}),
           ("   ·   서버 안 ", {}), ("180s+", {"color": SIGNAL, "bold": True}), (" 미완", {})],
          [("별도 프로세스로 분리 → 48.2s · 그래도 게임에는 느려 채택 보류", {"color": MUTED})]]),
        ("05", "재현 먼저, 수정은 그 다음",
         "배경 오염도 host 정지도 하니스로 먼저 재현한 뒤 고쳤다. 고치고 나서 통과하는 것보다, "
         "일부러 되돌려 실패하는 것을 확인하는 편이 확실하다.",
         [[("검증 하니스 ", {}), ("13종", {"color": ACCENT, "bold": True}),
           ("  —  로직 7 · 브라우저 6 (렌더 픽셀 검사 2종 포함)", {})]]),
    ]:
        box(s, M, y, W - 2 * M, Emu(9525), fill=RULE)
        text(s, M, Emu(y + inch(0.14)), inch(0.4), inch(0.25),
             [(no, {"font": MONO, "size": 11, "color": ACCENT, "bold": True})])
        tx = Emu(M + inch(0.45))
        text(s, tx, Emu(y + inch(0.12)), inch(10.3), inch(0.26), [(t, {"size": 14, "bold": True})])
        text(s, tx, Emu(y + inch(0.44)), inch(10.3), inch(0.62), [d], size=11, color=MUTED, spacing=1.28)
        evidence(s, tx, Emu(y + inch(1.12)), inch(10.2), ev, h=inch(0.46))
        y = Emu(y + inch(1.74))

    # ── 7. 한계 ───────────────────────────────────────────────────
    s = slide_base(prs, "LIMITS", "한계", None)
    cw, ch = inch(5.35), inch(1.62)
    for i, (tag, t, b) in enumerate([
        ("구조", "뒤통수·머리카락 없음",
         "모델이 얼굴 앞면만 담는다. 얼굴을 등록하면 아바타 머리를 통째로 숨기고 얼굴만 남긴다 — "
         "정면 게임이라 성립하는 타협이다."),
        ("촬영", "조명이 텍스처에 구워진다",
         "사진의 그림자가 그대로 3D에 남는다. 창을 등지면 얼굴이 실루엣이 된다. "
         "반사율과 조명을 분리하지 않았다."),
        ("깊이", "옆 얼굴은 추정치",
         "정면 한 장이라 코·턱 옆선의 깊이가 관측이 아니라 추정이다. 옆모습 2장을 더 찍는 "
         "방식은 조명·각도가 튀어 오히려 나빠졌다."),
        ("표정", "규칙 기반 오프셋",
         "피격·지침 표정은 미리 정의한 정점 변위다. 실시간 표정 추적이 아니라 HP에 반응하는 연출이다."),
    ]):
        x = M if i % 2 == 0 else Emu(M + cw + inch(0.20))
        y = inch(2.00) if i < 2 else inch(3.80)
        card(s, x, y, cw, ch, tag, SIGNAL, t, b)

    # ── 8. 다음 단계 ──────────────────────────────────────────────
    s = slide_base(prs, "NEXT", "다음에 할 것",
                   "뒤통수를 제대로 만들려면 얼굴만 담긴 모델을 벗어나야 한다. 각각 무엇을 포기하는지가 다르다.")
    hdr_y = inch(2.05)
    cols = [inch(2.55), inch(3.55), inch(3.35), inch(1.45)]
    xs, acc = [], M
    for c in cols:
        xs.append(acc)
        acc = Emu(acc + c)
    box(s, M, hdr_y, W - 2 * M, inch(0.34), fill=WASH, line=RULE)
    for x, c, h in zip(xs, cols, ["방향", "얻는 것", "포기하는 것", "비용"]):
        text(s, Emu(x + inch(0.14)), Emu(hdr_y + inch(0.09)), Emu(c - inch(0.2)), inch(0.2),
             [(h, {"font": MONO, "size": 9, "color": MUTED, "bold": True})])
    rows = [
        ("FLAME (MICA/DECA)", "두개골·귀·목이 포함된 통계 모델. 머리 형태가 구조적으로 해결된다",
         "머리카락이 없다 — 여전히 색으로 칠해야 한다. 연구용 라이선스", "GPU ~0.3s", True),
        ("TripoSR / TRELLIS", "머리카락까지 포함한 전체 메시를 사진 한 장에서 생성",
         "게임에 쓰기엔 느리다. 별도 워커 프로세스 필수", "2060 16~48s", False),
        ("조명 정규화", "역광·그림자에 강해진다. 촬영 성공률이 올라간다",
         "정확한 albedo 분리는 그 자체로 별도 연구 주제", "전처리", False),
        ("홍채 468→478", "눈동자 위치가 정확해진다. 촬영 순간에만 켜면 된다",
         "거의 없음 — 손이 가장 적게 드는 개선", "옵션 한 줄", False),
    ]
    ry = Emu(hdr_y + inch(0.34))
    for name, gain, cost, price, pick in rows:
        rh = inch(0.72)
        box(s, M, ry, W - 2 * M, rh, fill=CARD, line=RULE)
        text(s, Emu(xs[0] + inch(0.14)), Emu(ry + inch(0.12)), Emu(cols[0] - inch(0.24)), inch(0.24),
             [(name, {"size": 11.5, "bold": True})])
        if pick:
            text(s, Emu(xs[0] + inch(0.14)), Emu(ry + inch(0.38)), Emu(cols[0] - inch(0.24)), inch(0.2),
                 [("가장 유력", {"font": MONO, "size": 9, "color": ACCENT, "bold": True})])
        text(s, Emu(xs[1] + inch(0.14)), Emu(ry + inch(0.12)), Emu(cols[1] - inch(0.24)), inch(0.5),
             [gain], size=10.5, color=MUTED, spacing=1.2)
        text(s, Emu(xs[2] + inch(0.14)), Emu(ry + inch(0.12)), Emu(cols[2] - inch(0.24)), inch(0.5),
             [cost], size=10.5, color=MUTED, spacing=1.2)
        text(s, Emu(xs[3] + inch(0.14)), Emu(ry + inch(0.12)), Emu(cols[3] - inch(0.24)), inch(0.24),
             [(price, {"font": MONO, "size": 10.5, "color": INK})])
        ry = Emu(ry + rh)

    card(s, M, inch(5.62), inch(10.9), inch(0.78), "방법론", ACCENT,
         "픽셀 검사를 기본값으로",
         "렌더해서 색을 세는 하니스 2종이 좌표 검사가 놓친 버그를 전부 잡았다. 새 기능마다 먼저 만들 것.")

    prs.save(OUT)
    return OUT


if __name__ == "__main__":
    try:
        sys.stdout.reconfigure(encoding="utf-8", errors="replace")
    except Exception:
        pass
    p = build()
    print(f"저장: {p}  ({p.stat().st_size / 1024:.0f}KB)")
